from __future__ import annotations

import io
import zipfile
from pathlib import Path
from types import ModuleType
from typing import Any

import pytest
from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.presentation import Presentation as PresentationType

from conftest import generate_deck


@pytest.mark.parametrize("accent", ["FDE047", "111827", "777777", "0EA5E9"])
def test_custom_palette_is_derived_and_contrast_safe(
    powerpointgen: ModuleType, accent: str, tmp_path: Path
) -> None:
    options = powerpointgen.parse_options(
        {
            "slides": [{"layout": "title"}],
            "output_path": str(tmp_path / "x.pptx"),
            "accent_color": accent,
        }
    )
    palette = options.palette
    assert palette.accent == accent
    assert palette.accent_soft != powerpointgen.THEMES["executive"].accent_light
    assert powerpointgen._contrast_ratio(palette.on_accent, palette.accent) >= 4.5
    assert powerpointgen._contrast_ratio(palette.accent_text, "FFFFFF") >= 4.5
    assert powerpointgen._contrast_ratio(palette.link, "FFFFFF") >= 4.5
    assert powerpointgen._contrast_ratio(palette.link_on_accent, palette.accent) >= 4.5


def test_font_resolution_is_deterministic_and_honest(powerpointgen: ModuleType) -> None:
    catalog = {
        "liberation sans": ("Liberation Sans", "/fonts/LiberationSans.ttf"),
        "dejavu sans mono": ("DejaVu Sans Mono", "/fonts/DejaVuSansMono.ttf"),
    }
    resolved = powerpointgen._resolve_font("body", "Missing Primary", catalog)
    assert resolved.requested == "Missing Primary"
    assert resolved.resolved == "Liberation Sans"
    assert resolved.substituted is True
    assert resolved.verified is True
    assert resolved.metrics_path == "/fonts/LiberationSans.ttf"

    unresolved = powerpointgen._resolve_font("body", "Missing Primary", {})
    assert unresolved.resolved == "Missing Primary"
    assert unresolved.substituted is False
    assert unresolved.verified is False
    assert unresolved.metrics_path is None


def test_font_discovery_includes_user_windows_fonts_and_all_collection_faces(
    tmp_path: Path, powerpointgen: ModuleType, monkeypatch: pytest.MonkeyPatch
) -> None:
    local_app_data = tmp_path / "LocalAppData"
    user_fonts = local_app_data / "Microsoft" / "Windows" / "Fonts"
    user_fonts.mkdir(parents=True)
    monkeypatch.setenv("LOCALAPPDATA", str(local_app_data))
    assert user_fonts in powerpointgen._font_directories()

    collection = tmp_path / "collection.ttc"
    collection.write_bytes(b"fake")
    monkeypatch.setattr(powerpointgen, "_font_directories", lambda: [tmp_path])

    class FakeFont:
        def __init__(self, family: str, style: str) -> None:
            self.family = family
            self.style = style

        def getname(self) -> tuple[str, str]:
            return self.family, self.style

    faces = {
        0: ("Family A", "Bold"),
        1: ("Family A", "Regular"),
        2: ("Family B", "Regular"),
        3: ("Family A", "Italic"),
        4: ("Family A", "Bold Italic"),
    }

    def fake_truetype(_path: str, _size: int, index: int = 0) -> FakeFont:
        if index not in faces:
            raise OSError("no more faces")
        return FakeFont(*faces[index])

    monkeypatch.setattr(powerpointgen.ImageFont, "truetype", fake_truetype)
    powerpointgen._font_catalog.cache_clear()
    powerpointgen._font_style_catalog.cache_clear()
    try:
        catalog = powerpointgen._font_catalog()
        assert catalog["family a"] == ("Family A", str(collection), 1)
        assert catalog["family b"] == ("Family B", str(collection), 2)
        styles = powerpointgen._font_style_catalog()["family a"]
        assert styles == {
            "bold": (str(collection), 0),
            "regular": (str(collection), 1),
            "italic": (str(collection), 3),
            "bold_italic": (str(collection), 4),
        }
        resolved = powerpointgen._resolve_font("body", "Family A")
        assert resolved.metrics_face_index == 1
        assert resolved.metrics_styles == styles
    finally:
        powerpointgen._font_catalog.cache_clear()
        powerpointgen._font_style_catalog.cache_clear()


def test_generation_reports_font_roles_and_palette(
    powerpointgen: ModuleType, tmp_path: Path
) -> None:
    result, _ = generate_deck(
        powerpointgen,
        tmp_path,
        slides=[{"layout": "title", "title": "Fonts"}],
        accent_color="FDE047",
    )
    assert {record["role"] for record in result["fonts"]} == {"heading", "body", "mono"}
    assert all("verified" in record for record in result["fonts"])
    assert all(
        set(record["metrics_styles"]) == {"regular", "bold", "italic", "bold_italic"}
        for record in result["fonts"]
        if record["verified"]
    )
    assert result["resolved_palette"]["accent"] == "FDE047"


def test_bullet_glyph_uses_resolved_body_font(powerpointgen: ModuleType, tmp_path: Path) -> None:
    result, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        slides=[{"layout": "content", "bullets": ["Bullet"]}],
    )
    body_font = next(record["resolved"] for record in result["fonts"] if record["role"] == "body")
    bullet_fonts = presentation.slides[0]._element.xpath(".//a:buFont/@typeface")
    assert bullet_fonts == [body_font]


def test_light_accent_section_table_and_closing_use_on_accent(
    powerpointgen: ModuleType, tmp_path: Path
) -> None:
    result, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        slides=[
            {"layout": "section", "title": "Section"},
            {
                "layout": "table",
                "title": "Table",
                "table": {"headers": ["Header"], "rows": [["Value"]]},
            },
            {"layout": "closing", "title": "Closing", "subtitle": "Subtitle"},
        ],
        accent_color="FDE047",
    )
    on_accent = result["resolved_palette"]["on_accent"]
    expected = ["Section", "Header", "Closing"]
    for slide, text in zip(presentation.slides, expected, strict=True):
        runs: list[Any] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                runs.extend(
                    run for paragraph in shape.text_frame.paragraphs for run in paragraph.runs
                )
            if getattr(shape, "has_table", False):
                runs.extend(
                    run
                    for row in shape.table.rows
                    for cell in row.cells
                    for paragraph in cell.text_frame.paragraphs
                    for run in paragraph.runs
                )
        run = next(candidate for candidate in runs if candidate.text == text)
        assert str(run.font.color.rgb) == on_accent


def _picture_shapes(presentation: PresentationType) -> list[Any]:
    return [
        shape
        for slide in presentation.slides
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]


@pytest.mark.parametrize("size", [(120, 600), (800, 120)])
def test_contain_preserves_image_aspect_ratio(
    tmp_path: Path, powerpointgen: ModuleType, size: tuple[int, int]
) -> None:
    image_path = tmp_path / f"{size[0]}x{size[1]}.png"
    Image.new("RGB", size, "red").save(image_path)
    _, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        name=f"contain-{size[0]}",
        slides=[{"layout": "image", "image_path": image_path.name}],
    )
    picture = _picture_shapes(presentation)[0]
    assert picture.width / picture.height == pytest.approx(size[0] / size[1], rel=0.01)


def test_exif_orientation_is_normalized_in_embedded_image(
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    image_path = tmp_path / "oriented.jpg"
    image = Image.new("RGB", (100, 200), "blue")
    exif = image.getexif()
    exif[274] = 6
    image.save(image_path, exif=exif)

    result, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        slides=[{"layout": "image", "image_path": image_path.name}],
    )
    picture = _picture_shapes(presentation)[0]
    assert picture.width / picture.height == pytest.approx(2.0, rel=0.01)

    with zipfile.ZipFile(result["path"]) as archive:
        media_name = next(name for name in archive.namelist() if name.startswith("ppt/media/"))
        with Image.open(io.BytesIO(archive.read(media_name))) as embedded:
            assert embedded.size == (200, 100)
            assert embedded.getexif().get(274) is None


def test_low_resolution_image_warns(tmp_path: Path, powerpointgen: ModuleType) -> None:
    image_path = tmp_path / "tiny.png"
    Image.new("RGB", (24, 24), "green").save(image_path)
    result, _ = generate_deck(
        powerpointgen,
        tmp_path,
        slides=[{"layout": "image", "image_path": image_path.name}],
    )
    assert any("low effective" in warning for warning in result["warnings"])


def test_corrupt_existing_image_fails_and_preserves_target(
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    image_path = tmp_path / "corrupt.png"
    image_path.write_bytes(b"not an image")
    target = tmp_path / "existing.pptx"
    target.write_bytes(b"known-good-placeholder")
    with pytest.raises(ValueError, match="corrupt image"):
        powerpointgen.generate(
            {
                "slides": [{"layout": "image", "image_path": image_path.name}],
                "project_root": str(tmp_path),
                "output_path": str(target),
            }
        )
    assert target.read_bytes() == b"known-good-placeholder"
