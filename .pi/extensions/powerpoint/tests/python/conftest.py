from __future__ import annotations

import importlib.util
import sys
from pathlib import Path
from types import ModuleType
from typing import Any

import pytest
from pptx import Presentation
from pptx.presentation import Presentation as PresentationType

GENERATOR = Path(__file__).parents[2] / "generate_pptx.py"


@pytest.fixture(scope="session")
def powerpointgen() -> ModuleType:
    spec = importlib.util.spec_from_file_location("penny_powerpoint_generate_pptx", GENERATOR)
    assert spec is not None and spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


def generate_deck(
    powerpointgen: ModuleType,
    tmp_path: Path,
    *,
    name: str = "deck",
    **spec: Any,
) -> tuple[dict[str, Any], PresentationType]:
    output = tmp_path / f"{name}.pptx"
    payload = {
        "output_path": str(output),
        "project_root": str(tmp_path),
        **spec,
    }
    result = powerpointgen.generate(payload)
    return result, Presentation(str(output))


def slide_texts(presentation: PresentationType) -> list[str]:
    texts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                texts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    texts.extend(cell.text for cell in row.cells)
    return texts
