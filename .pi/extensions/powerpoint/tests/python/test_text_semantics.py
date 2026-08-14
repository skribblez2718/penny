from __future__ import annotations

from pathlib import Path
from types import ModuleType
from typing import Any

import pytest
from pptx.presentation import Presentation as PresentationType

from conftest import generate_deck

MARKDOWN = (
    "## Semantics\n\n"
    "Soft A\nSoft B\n\n"
    "Hard A  \nHard B\n\n"
    "HTML A<br>HTML B\n\n"
    "First paragraph.\n\n"
    "Second paragraph with [linked **words**](https://example.com/path).\n"
)


def _body_shape(presentation: PresentationType) -> Any:
    return next(
        shape
        for shape in presentation.slides[0].shapes
        if getattr(shape, "has_text_frame", False) and "Soft A" in shape.text
    )


def test_preserve_mode_keeps_soft_hard_html_breaks_and_paragraphs(
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    _, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        markdown=MARKDOWN,
        line_break_mode="preserve",
    )

    shape = _body_shape(presentation)
    paragraphs = shape.text_frame.paragraphs
    assert len(paragraphs) == 5
    assert paragraphs[0].text == "Soft A\vSoft B"
    assert paragraphs[1].text == "Hard A\vHard B"
    assert paragraphs[2].text == "HTML A\vHTML B"
    assert shape._element.xpath(".//a:br")


def test_commonmark_only_folds_soft_break(tmp_path: Path, powerpointgen: ModuleType) -> None:
    _, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        markdown=MARKDOWN,
        line_break_mode="commonmark",
    )

    paragraphs = _body_shape(presentation).text_frame.paragraphs
    assert paragraphs[0].text == "Soft A Soft B"
    assert paragraphs[1].text == "Hard A\vHard B"
    assert paragraphs[2].text == "HTML A\vHTML B"


def test_hyperlink_survives_emphasis_and_heading(tmp_path: Path, powerpointgen: ModuleType) -> None:
    _, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        markdown="## [Linked heading](https://example.com/heading)\n\n"
        "A [linked **phrase**](https://example.com/body).",
    )

    addresses = [
        run.hyperlink.address
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.hyperlink.address
    ]
    assert "https://example.com/heading" in addresses
    assert addresses.count("https://example.com/body") >= 2
    linked_runs = [
        run
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.hyperlink.address
    ]
    assert all(run.font.underline for run in linked_runs)


@pytest.mark.parametrize(
    "markdown",
    [
        "## Body\n\nText before ![icon](missing.png) text after.",
        "## Heading\n\n- Text ![icon](missing.png)",
        "## Table\n\n| A |\n|---|\n| Text ![icon](missing.png) |",
        "## [Heading ![icon](missing.png)](https://example.com)",
    ],
)
def test_mixed_inline_images_fail_without_publishing(
    tmp_path: Path, powerpointgen: ModuleType, markdown: str
) -> None:
    output = tmp_path / "must-not-exist.pptx"
    with pytest.raises(ValueError, match="inline image"):
        powerpointgen.generate(
            {
                "markdown": markdown,
                "output_path": str(output),
                "project_root": str(tmp_path),
            }
        )
    assert not output.exists()


@pytest.mark.parametrize(
    "markdown",
    [
        "## Raw HTML\n\nBefore <img src='missing.png' alt='raw'> after.",
        "## Heading <img src='missing.png' alt='raw'>",
        "## List\n\n- Before <img src='missing.png' alt='raw'> after",
        "## Table\n\n| A |\n|---|\n| Before <img src='missing.png' alt='raw'> after |",
    ],
)
def test_raw_html_images_fail_closed_and_preserve_target(
    tmp_path: Path, powerpointgen: ModuleType, markdown: str
) -> None:
    output = tmp_path / "existing.pptx"
    output.write_bytes(b"baseline")
    with pytest.raises(ValueError, match="inline image"):
        powerpointgen.generate(
            {
                "markdown": markdown,
                "output_path": str(output),
                "project_root": str(tmp_path),
            }
        )
    assert output.read_bytes() == b"baseline"


def test_table_intro_and_two_column_keep_paragraph_boundaries(
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    _, markdown_presentation = generate_deck(
        powerpointgen,
        tmp_path,
        name="table-paragraphs",
        markdown=("## Table\n\nFIRST-PARA\n\nSECOND-PARA\n\n" "| Header |\n|---|\n| Value |"),
    )
    intro = next(
        shape
        for shape in markdown_presentation.slides[0].shapes
        if getattr(shape, "has_text_frame", False) and "FIRST-PARA" in shape.text
    )
    assert [paragraph.text for paragraph in intro.text_frame.paragraphs] == [
        "FIRST-PARA",
        "SECOND-PARA",
    ]

    _, structured_presentation = generate_deck(
        powerpointgen,
        tmp_path,
        name="column-paragraphs",
        slides=[
            {
                "layout": "two_column",
                "left": {"body": "FIRST-COLUMN\n\nSECOND-COLUMN"},
                "right": {"body": "Right"},
            }
        ],
    )
    column_body = next(
        shape
        for shape in structured_presentation.slides[0].shapes
        if getattr(shape, "has_text_frame", False) and "FIRST-COLUMN" in shape.text
    )
    assert [paragraph.text for paragraph in column_body.text_frame.paragraphs] == [
        "FIRST-COLUMN",
        "SECOND-COLUMN",
    ]


def test_break_policy_applies_in_bullet_and_table_cell(
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    _, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        slides=[
            {"layout": "content", "bullets": ["Bullet A\nBullet B"]},
            {
                "layout": "table",
                "table": {
                    "headers": ["Head A<br>Head B"],
                    "rows": [["Cell A\nCell B"]],
                },
            },
        ],
    )
    assert sum(len(slide._element.xpath(".//a:br")) for slide in presentation.slides) == 3


def test_link_with_break_keeps_target_on_runs_on_both_sides(
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    _, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        markdown="## Link\n\n[Before  \nAfter](https://example.com/break)",
    )
    linked = [
        run
        for shape in presentation.slides[0].shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.hyperlink.address
    ]
    assert [run.text for run in linked] == ["Before", "After"]
    assert all(run.hyperlink.address == "https://example.com/break" for run in linked)
    assert presentation.slides[0]._element.xpath(".//a:br")


def test_image_only_paragraph_remains_supported(tmp_path: Path, powerpointgen: ModuleType) -> None:
    result, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        markdown="![diagram](missing.png)",
    )
    assert result["layouts_used"] == {"image": 1}
    assert "[image unavailable: diagram]" in "\n".join(
        shape.text
        for shape in presentation.slides[0].shapes
        if getattr(shape, "has_text_frame", False)
    )
