from __future__ import annotations

import warnings
import zipfile
from pathlib import Path
from types import ModuleType
from typing import Any, cast

import pytest
from pptx import Presentation

from conftest import generate_deck

REQUIRED_PARTS = {
    "[Content_Types].xml",
    "_rels/.rels",
    "docProps/core.xml",
    "docProps/app.xml",
    "ppt/presentation.xml",
    "ppt/_rels/presentation.xml.rels",
}


def test_generated_package_passes_crc_xml_reopen_and_slide_count(
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    result, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        slides=[
            {"layout": "title", "title": "Validated"},
            {"layout": "content", "title": "Body", "body": "Hello."},
        ],
    )
    validation = result["validation"]
    assert validation["package_valid"] is True
    assert validation["reopen_valid"] is True
    assert validation["slide_count"] == result["slide_count"] == len(presentation.slides)
    assert validation["openxml_schema_validation"] == "not_performed"
    with zipfile.ZipFile(result["path"]) as archive:
        assert REQUIRED_PARTS <= set(archive.namelist())
        assert archive.testzip() is None


def test_validation_rejects_duplicate_part_names(tmp_path: Path, powerpointgen: ModuleType) -> None:
    invalid = tmp_path / "duplicate.pptx"
    with warnings.catch_warnings():
        warnings.simplefilter("ignore", UserWarning)
        with zipfile.ZipFile(invalid, "w") as archive:
            archive.writestr("ppt/presentation.xml", "<p:presentation xmlns:p='x'/>")
            archive.writestr("ppt/presentation.xml", "<p:presentation xmlns:p='x'/>")
    with pytest.raises(ValueError, match="duplicate"):
        powerpointgen._validate_pptx(invalid)


@pytest.mark.parametrize(
    "payload",
    [
        b'<!DOCTYPE Properties [<!ENTITY injected "expanded">]><Properties>&injected;</Properties>',
        b'<!DOCTYPE Properties [<!ENTITY external SYSTEM "file:///etc/passwd">]><Properties>&external;</Properties>',
        b'<!DOCTYPE Properties [<!ENTITY % external SYSTEM "http://127.0.0.1:9/evil.dtd">%external;]><Properties/>',
    ],
)
def test_validation_rejects_dtd_and_entity_bearing_xml(
    tmp_path: Path, powerpointgen: ModuleType, payload: bytes
) -> None:
    source = tmp_path / "source.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(str(source))
    adversarial = tmp_path / "entity-bearing.pptx"
    with zipfile.ZipFile(source) as original, zipfile.ZipFile(adversarial, "w") as output:
        for info in original.infolist():
            content = payload if info.filename == "docProps/app.xml" else original.read(info)
            output.writestr(info, content)

    with pytest.raises(ValueError, match="forbidden DTD/entity|XML parse failure"):
        powerpointgen._validate_pptx(adversarial)


def test_package_xml_parser_disables_resolution_and_network(
    powerpointgen: ModuleType, monkeypatch: pytest.MonkeyPatch
) -> None:
    captured: dict[str, Any] = {}
    original = powerpointgen.etree.XMLParser

    def capture_parser(**kwargs: Any) -> Any:
        captured.update(kwargs)
        return original(**kwargs)

    monkeypatch.setattr(powerpointgen.etree, "XMLParser", capture_parser)
    powerpointgen._parse_package_xml(b"<root/>", "safe.xml")
    assert captured == {
        "resolve_entities": False,
        "no_network": True,
        "load_dtd": False,
        "recover": False,
        "huge_tree": False,
    }


def test_validation_rejects_missing_part_and_malformed_xml(
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    missing = tmp_path / "missing.pptx"
    with zipfile.ZipFile(missing, "w") as archive:
        archive.writestr("[Content_Types].xml", "<Types/>")
    with pytest.raises(ValueError, match="missing required parts"):
        powerpointgen._validate_pptx(missing)

    malformed = tmp_path / "malformed.pptx"
    source = Presentation()
    source.slides.add_slide(source.slide_layouts[6])
    source.save(str(malformed))
    rewritten = tmp_path / "rewritten.pptx"
    with zipfile.ZipFile(malformed) as original, zipfile.ZipFile(rewritten, "w") as output:
        for info in original.infolist():
            payload = b"<broken" if info.filename == "ppt/presentation.xml" else original.read(info)
            output.writestr(info, payload)
    with pytest.raises(ValueError, match="XML|xml|parse"):
        powerpointgen._validate_pptx(rewritten)


def test_atomic_validation_failure_preserves_existing_target_and_cleans_stage(
    tmp_path: Path, powerpointgen: ModuleType, monkeypatch: pytest.MonkeyPatch
) -> None:
    target = tmp_path / "existing.pptx"
    target.write_bytes(b"baseline")
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])

    def reject(_path: Path, expected_slide_count: int | None = None) -> None:
        del expected_slide_count
        raise ValueError("forced validation failure")

    monkeypatch.setattr(powerpointgen, "_validate_pptx", reject)
    with pytest.raises(ValueError, match="forced validation failure"):
        powerpointgen._atomic_save(presentation, target, staging_path=None)
    assert target.read_bytes() == b"baseline"
    assert not list(tmp_path.glob(".*.tmp.pptx"))


def test_atomic_success_replaces_target_and_uses_same_directory_stage(
    tmp_path: Path, powerpointgen: ModuleType, monkeypatch: pytest.MonkeyPatch
) -> None:
    target = tmp_path / "existing.pptx"
    target.write_bytes(b"baseline")
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    observed: list[Path] = []
    original = powerpointgen._validate_pptx

    def observe(
        path: Path,
        expected_slide_count: int | None = None,
    ) -> dict[str, Any]:
        observed.append(Path(path))
        return cast(
            dict[str, Any],
            original(path, expected_slide_count=expected_slide_count),
        )

    monkeypatch.setattr(powerpointgen, "_validate_pptx", observe)
    resolved, validation = powerpointgen._atomic_save(presentation, target, staging_path=None)
    assert Path(resolved) == target
    assert validation["slide_count"] == 1
    assert observed[0].parent == target.parent
    assert target.read_bytes().startswith(b"PK")
    assert not observed[0].exists()


def test_explicit_staging_must_be_same_directory(tmp_path: Path, powerpointgen: ModuleType) -> None:
    target = tmp_path / "target" / "deck.pptx"
    staging = tmp_path / "elsewhere" / ".deck.tmp.pptx"
    staging.parent.mkdir()
    presentation = Presentation()
    with pytest.raises(ValueError, match="same directory"):
        powerpointgen._atomic_save(presentation, target, staging_path=staging)


def test_all_public_layouts_and_themes_reopen(tmp_path: Path, powerpointgen: ModuleType) -> None:
    image_path = tmp_path / "image.png"
    from PIL import Image

    Image.new("RGB", (200, 100), "purple").save(image_path)
    layouts = [
        {"layout": "title", "title": "Title"},
        {"layout": "section", "title": "Section"},
        {
            "layout": "content",
            "title": "Content",
            "body": "Body",
            "bullets": ["Bullet"],
            "notes": "Speaker note",
        },
        {
            "layout": "two_column",
            "title": "Columns",
            "left": {"heading": "Left", "body": "Body"},
            "right": {"heading": "Right", "bullets": ["Bullet"]},
        },
        {"layout": "table", "title": "Table", "table": {"headers": ["A"], "rows": [["B"]]}},
        {"layout": "quote", "quote": "Quote", "attribution": "Person"},
        {"layout": "image", "image_path": image_path.name, "caption": "Image"},
        {"layout": "closing", "title": "Close", "subtitle": "Done"},
    ]
    for theme in powerpointgen.THEMES:
        result, presentation = generate_deck(
            powerpointgen,
            tmp_path,
            name=f"theme-{theme}",
            slides=layouts,
            theme=theme,
            footer_text="Internal",
        )
        assert result["slide_count"] == len(layouts)
        assert len(presentation.slides) == len(layouts)
        assert "Speaker note" in presentation.slides[2].notes_slide.notes_text_frame.text
        assert any(
            shape.text == "Internal"
            for shape in presentation.slides[2].shapes
            if getattr(shape, "has_text_frame", False)
        )
