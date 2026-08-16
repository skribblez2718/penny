from __future__ import annotations

import random
import zipfile
from pathlib import Path
from types import ModuleType
from typing import Any

import pytest
from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE

from conftest import generate_deck, slide_texts


def _assets(tmp_path: Path) -> tuple[Path, Path]:
    media = tmp_path / "media.jpg"
    mark = tmp_path / "mark.png"
    Image.new("RGB", (1200, 600), "orange").save(media)
    Image.new("RGBA", (240, 120), (20, 200, 120, 180)).save(mark)
    return media, mark


def _shape_named(slide: Any, name: str) -> Any:
    return next(shape for shape in slide.shapes if shape.name == name)


def _base_design(media: Path, mark: Path) -> dict[str, Any]:
    return {
        "palette": {
            "canvas": "0B1020",
            "surface": "151F32",
            "accent": "38BDF8",
            "text": "222222",
            "muted_text": "333333",
        },
        "background": {
            "path": media.name,
            "fit": "crop",
            "focal_point": {"x": 0.9, "y": 0.5},
            "overlay": {"opacity": 0.4},
        },
        "mark": {
            "path": mark.name,
            "x": 0.80,
            "y": 0.10,
            "width": 0.10,
            "height": 0.20,
        },
    }


@pytest.mark.parametrize(
    "spec,match",
    [
        ({"design": {"unknown": 1}}, "unknown keys"),
        ({"design": {"palette": {"canvas": "12345"}}}, "strict 6-digit"),
        (
            {"design": {"background": {"path": "a.png", "unknown": True}}},
            "unknown keys",
        ),
        (
            {
                "design": {
                    "background": {
                        "path": "a.png",
                        "focal_point": {"x": 0.5, "y": 2},
                    }
                }
            },
            "must be in",
        ),
        (
            {
                "design": {
                    "mark": {
                        "path": "a.png",
                        "x": 0.9,
                        "y": 0,
                        "width": 0.2,
                        "height": 0.1,
                    }
                }
            },
            "inside normalized",
        ),
        (
            {"design": {"background": {"path": "a.png", "overlay": {}}}},
            "opacity is required",
        ),
        (
            {"design": {"background": {"path": "a.png", "fit": None}}},
            "must be 'crop' or 'contain'",
        ),
        (
            {"design": {"background": {"path": "a.png", "focal_point": None}}},
            "must be an object",
        ),
        (
            {"design": {"background": {"path": "a.png", "overlay": None}}},
            "must be an object",
        ),
        (
            {
                "design": {
                    "mark": {
                        "path": "a.png",
                        "x": 0,
                        "y": 0,
                        "width": 0.1,
                        "height": 0.1,
                        "fit": None,
                    }
                }
            },
            "must be 'crop' or 'contain'",
        ),
    ],
)
def test_new_nested_contract_is_strict(
    tmp_path: Path,
    powerpointgen: ModuleType,
    spec: dict[str, Any],
    match: str,
) -> None:
    with pytest.raises(ValueError, match=match):
        powerpointgen.parse_options(
            {
                "slides": [{"layout": "title"}],
                "output_path": str(tmp_path / "x.pptx"),
                **spec,
            }
        )


def test_defaults_conflict_and_composed_layout_validation(
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    options = powerpointgen.parse_options(
        {
            "slides": [{"layout": "title"}],
            "output_path": str(tmp_path / "x.pptx"),
            "design": {
                "background": {"path": "a.png"},
                "mark": {
                    "path": "b.png",
                    "x": 0,
                    "y": 0,
                    "width": 0.1,
                    "height": 0.1,
                },
            },
        }
    )
    assert options.design.background is not None
    assert options.design.background.fit == "crop"
    assert options.design.mark is not None
    assert options.design.mark.fit == "contain"

    with pytest.raises(ValueError, match="cannot be combined"):
        powerpointgen.parse_options(
            {
                "slides": [{"layout": "title"}],
                "output_path": str(tmp_path / "conflict.pptx"),
                "accent_color": "112233",
                "design": {"palette": {"accent": "445566"}},
            }
        )
    with pytest.raises(ValueError, match="requires media"):
        powerpointgen.normalize_slide({"layout": "image_left", "title": "Missing"})
    with pytest.raises(ValueError, match="incompatible"):
        powerpointgen.normalize_slide(
            {
                "layout": "image_right",
                "media": {"path": "a.png"},
                "table": {"headers": [], "rows": []},
            }
        )
    slide = powerpointgen.normalize_slide({"layout": "full_bleed", "media": {"path": "a.png"}})
    assert slide["_media"].fit == "crop"

    with pytest.raises(ValueError, match="non-empty array"):
        powerpointgen.parse_options(
            {
                "slides": [{"layout": "title"}],
                "output_path": str(tmp_path / "empty-roots.pptx"),
                "allowed_image_roots": [],
            }
        )


def test_palette_merge_inherit_null_reset_and_no_leak(
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    media, mark = _assets(tmp_path)
    result, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        design=_base_design(media, mark),
        slides=[
            {
                "layout": "content",
                "title": "Inherited",
                "body": "Body",
                "design": {"palette": {"surface": "202A44"}},
            },
            {
                "layout": "content",
                "title": "Removed",
                "body": "Body",
                "design": {"background": None, "mark": None},
            },
            {"layout": "content", "title": "Reset", "body": "Body"},
        ],
    )
    records = result["resolved_design"]["slides"]
    assert records[0]["palette"]["surface"] == "202A44"
    assert records[0]["palette"]["canvas"] == "0B1020"
    assert records[0]["background"]["present"] is True
    assert records[0]["mark"]["present"] is True
    assert records[1]["background"]["present"] is False
    assert records[1]["mark"]["present"] is False
    assert records[2]["palette"]["surface"] == "151F32"
    assert records[2]["background"]["present"] is True
    assert records[2]["mark"]["present"] is True
    assert all(_shape_named(slide, "Penny Canvas") for slide in presentation.slides)


def test_design_contrast_corrections_and_dark_semantic_surfaces(
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    result, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        design={
            "palette": {
                "canvas": "080B12",
                "surface": "101827",
                "accent": "FDE047",
                "text": "111111",
                "muted_text": "222222",
            }
        },
        slides=[
            {
                "layout": "content",
                "title": "Code",
                "body": "Readable body",
                "code": ["print('dark')"],
            },
            {
                "layout": "table",
                "title": "Table",
                "table": {"headers": ["Header"], "rows": [["One"], ["Two"]]},
            },
        ],
    )
    for record in result["resolved_design"]["slides"]:
        assert record["corrections"]
        assert all(role["ratio"] >= 4.5 for role in record["contrast_roles"].values())
    palette = result["resolved_design"]["deck_default"]["palette"]
    assert palette["surface"] != "FFFFFF"
    assert palette["surface_alt"] not in {"FFFFFF", "F5F7FA"}
    assert palette["code_background"] != "FFFFFF"

    slide_colors = [
        color
        for slide in presentation.slides
        for color in slide._element.xpath(".//a:solidFill/a:srgbClr/@val")
    ]
    assert "F5F7FA" not in slide_colors
    table = next(
        shape.table for shape in presentation.slides[1].shapes if getattr(shape, "has_table", False)
    )
    body_fills = [str(table.cell(row, 0).fill.fore_color.rgb) for row in (1, 2)]
    assert body_fills == [palette["surface"], palette["surface_alt"]]


def test_adversarial_and_randomized_palette_roles_always_meet_contrast_floor(  # noqa: C901
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    adversarial = {
        "canvas": "A75766",
        "surface": "A642DA",
        "accent": "EF59A2",
        "text": "806F3D",
        "muted_text": "E30339",
    }
    result, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        name="adversarial-palette",
        design={"palette": adversarial},
        slides=[
            {
                "layout": "content",
                "title": "Contrast sentinel",
                "body": "Body sentinel",
                "bullets": ["Bullet sentinel"],
                "code": ["print('code sentinel')"],
            },
            {
                "layout": "table",
                "title": "Table sentinel",
                "table": {"headers": ["Header"], "rows": [["One"], ["Two"]]},
            },
            {"layout": "section", "title": "Section sentinel"},
            {"layout": "closing", "title": "Closing sentinel"},
        ],
    )
    for record in result["resolved_design"]["slides"]:
        assert all(role["ratio"] >= 4.5 for role in record["contrast_roles"].values())

    content_surface = _shape_named(presentation.slides[0], "Penny Content Surface")
    surface = str(content_surface.fill.fore_color.rgb)
    assert surface == adversarial["surface"]
    for shape in presentation.slides[0].shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                color = run.font.color.rgb
                if color is not None:
                    assert powerpointgen._contrast_ratio(str(color), surface) >= 4.5

    table = next(
        shape.table for shape in presentation.slides[1].shapes if getattr(shape, "has_table", False)
    )
    for row in (1, 2):
        fill = str(table.cell(row, 0).fill.fore_color.rgb)
        for paragraph in table.cell(row, 0).text_frame.paragraphs:
            for run in paragraph.runs:
                assert run.font.color.rgb is not None
                assert powerpointgen._contrast_ratio(str(run.font.color.rgb), fill) >= 4.5

    legacy = powerpointgen._derive_palette(powerpointgen.THEMES["modern"])
    rng = random.Random(20260814)
    for _ in range(500):
        colors = [f"{rng.randrange(0x1000000):06X}" for _ in range(5)]
        palette, _ = powerpointgen._derive_design_palette(
            legacy,
            powerpointgen.PalettePatch(*colors),
        )
        pairs = [
            (palette.text, palette.surface),
            (palette.text, palette.surface_alt),
            (palette.text, palette.code_background),
            (palette.text_muted, palette.surface),
            (palette.text_muted, palette.surface_alt),
            (palette.text_muted, palette.code_background),
            (palette.accent_text, palette.surface),
            (palette.accent_text, palette.surface_alt),
            (palette.link, palette.surface_alt),
            (palette.code_text, palette.code_background),
            (palette.canvas_text, palette.background),
            (palette.canvas_text_muted, palette.background),
            (palette.on_accent, palette.accent),
            (palette.link_on_accent, palette.accent),
        ]
        assert all(
            powerpointgen._contrast_ratio(foreground, background) >= 4.5
            for foreground, background in pairs
        )


def test_accent_role_corrections_are_truthful_and_preserve_decorative_fill(
    powerpointgen: ModuleType,
) -> None:
    legacy = powerpointgen._derive_palette(powerpointgen.THEMES["modern"])
    palette, corrections = powerpointgen._derive_design_palette(
        legacy,
        powerpointgen.PalettePatch(
            canvas="050816",
            surface="101A2E",
            accent="0B2138",
            text="F8FAFC",
            muted_text="B6C2D5",
        ),
    )
    assert palette.accent == "0B2138"
    assert palette.accent_text != palette.accent
    assert {
        correction["actual_role"]
        for correction in corrections
        if correction["requested_role"] == "accent"
    } == {"accent_text", "link"}
    assert all(
        correction["decorative_value_preserved"] is True
        for correction in corrections
        if correction["requested_role"] == "accent"
    )


def test_background_crop_focal_overlay_order_and_mark_geometry(
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    media, mark = _assets(tmp_path)
    result, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        design=_base_design(media, mark),
        slides=[{"layout": "title", "title": "Layers"}],
    )
    slide = presentation.slides[0]
    names = [shape.name for shape in slide.shapes]
    assert names.index("Penny Canvas") < names.index("Penny Background")
    assert names.index("Penny Background") < names.index("Penny Background Overlay")
    assert names.index("Penny Background Overlay") < names.index("Penny Title Surface")
    assert names.index("Penny Title Surface") < names.index("Penny Mark")
    title_index = next(
        index for index, shape in enumerate(slide.shapes) if "Layers" in getattr(shape, "text", "")
    )
    assert names.index("Penny Mark") < title_index

    background = _shape_named(slide, "Penny Background")
    expected_visible = (powerpointgen.SLIDE_W / powerpointgen.SLIDE_H) / 2.0
    assert background.crop_left == pytest.approx(1.0 - expected_visible, abs=1e-4)
    assert background.crop_right == pytest.approx(0.0, abs=1e-4)
    overlay = _shape_named(slide, "Penny Background Overlay")
    assert overlay._element.xpath(".//a:alpha/@val") == ["40000"]

    mark_shape = _shape_named(slide, "Penny Mark")
    expected_x = 0.80 * powerpointgen.SLIDE_W
    expected_width = 0.10 * powerpointgen.SLIDE_W
    expected_height = expected_width / 2.0
    target_top = 0.10 * powerpointgen.SLIDE_H
    target_height = 0.20 * powerpointgen.SLIDE_H
    assert mark_shape.left / 914400 == pytest.approx(expected_x, abs=0.002)
    assert mark_shape.width / 914400 == pytest.approx(expected_width, abs=0.002)
    assert mark_shape.height / 914400 == pytest.approx(expected_height, abs=0.002)
    assert mark_shape.top / 914400 == pytest.approx(
        target_top + (target_height - expected_height) / 2,
        abs=0.002,
    )
    assert result["resolved_design"]["slides"][0]["background"]["fit"] == "crop"


def test_background_contain_preserves_aspect_and_letterboxes(
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    media, _ = _assets(tmp_path)
    _, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        design={"background": {"path": media.name, "fit": "contain"}},
        slides=[{"layout": "title", "title": "Contain"}],
    )
    picture = _shape_named(presentation.slides[0], "Penny Background")
    assert picture.width / picture.height == pytest.approx(2.0, rel=0.001)
    assert picture.width / 914400 == pytest.approx(powerpointgen.SLIDE_W, abs=0.002)
    assert picture.height / 914400 < powerpointgen.SLIDE_H
    assert picture.top / 914400 > 0


@pytest.mark.parametrize("layout", ["image_left", "image_right", "full_bleed"])
def test_composed_layouts_embed_media_and_keep_text_editable(
    tmp_path: Path, powerpointgen: ModuleType, layout: str
) -> None:
    media, _ = _assets(tmp_path)
    result, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        name=layout,
        design={"palette": {"canvas": "101827", "surface": "17243A"}},
        slides=[
            {
                "layout": layout,
                "title": "Editable title",
                "kicker": "Evidence",
                "body": "Editable body",
                "bullets": ["Editable bullet"],
                "caption": "Caption sentinel",
                "notes": "Speaker note",
                "media": {"path": media.name},
            }
        ],
    )
    slide = presentation.slides[0]
    media_shape = _shape_named(slide, "Penny Media")
    assert media_shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    joined = "\n".join(slide_texts(presentation))
    for sentinel in ("Editable title", "Editable body", "Editable bullet", "Caption sentinel"):
        assert joined.count(sentinel) == 1
    assert "Speaker note" in slide.notes_slide.notes_text_frame.text
    record = result["resolved_design"]["slides"][0]
    assert record["media"]["present"] is True
    assert record["media"]["fit"] == ("crop" if layout == "full_bleed" else "contain")

    if layout == "image_left":
        assert media_shape.left == 0
        assert media_shape.width / 914400 == pytest.approx(5.05, abs=0.002)
    elif layout == "image_right":
        assert media_shape.left / 914400 == pytest.approx(powerpointgen.SLIDE_W - 5.05, abs=0.002)
    else:
        panel = _shape_named(slide, "Penny Full Bleed Text Panel")
        assert (
            panel.left / 914400,
            panel.top / 914400,
            panel.width / 914400,
            panel.height / 914400,
        ) == pytest.approx(powerpointgen.FULL_BLEED_PANEL, abs=0.002)
        assert str(panel.fill.fore_color.rgb) == record["palette"]["surface"]

    with zipfile.ZipFile(result["path"]) as archive:
        rels = archive.read("ppt/slides/_rels/slide1.xml.rels").decode()
        assert 'TargetMode="External"' not in rels
        assert "../media/" in rels


def test_full_bleed_media_replaces_background_and_uses_real_focal_crop_overlay(
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    media, _ = _assets(tmp_path)
    _, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        design={"background": {"path": media.name}},
        slides=[
            {
                "layout": "full_bleed",
                "title": "Full",
                "media": {
                    "path": media.name,
                    "focal_point": {"x": 0.9, "y": 0.5},
                    "overlay": {"color": "112233", "opacity": 0.35},
                },
            }
        ],
    )
    slide = presentation.slides[0]
    names = [shape.name for shape in slide.shapes]
    assert "Penny Background" not in names
    assert names.index("Penny Canvas") < names.index("Penny Media")
    assert names.index("Penny Media") < names.index("Penny Media Overlay")
    assert names.index("Penny Media Overlay") < names.index("Penny Full Bleed Text Panel")
    picture = _shape_named(slide, "Penny Media")
    expected_visible = (powerpointgen.SLIDE_W / powerpointgen.SLIDE_H) / 2.0
    assert picture.crop_left == pytest.approx(1.0 - expected_visible, abs=1e-4)
    assert picture.crop_right == pytest.approx(0.0, abs=1e-4)
    overlay = _shape_named(slide, "Penny Media Overlay")
    assert overlay._element.xpath(".//a:alpha/@val") == ["35000"]


def test_composed_pagination_repeats_media_design_and_loses_no_sentinels(
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    media, mark = _assets(tmp_path)
    bullets = [f"SENTINEL-{index:02d} " + "detail " * 38 for index in range(12)]
    result, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        design={
            "palette": {"canvas": "090D17", "surface": "151F32"},
            "mark": {
                "path": mark.name,
                "x": 0.90,
                "y": 0.04,
                "width": 0.06,
                "height": 0.06,
            },
        },
        slides=[
            {
                "layout": "image_left",
                "title": "Continuation",
                "caption": "ONE-CAPTION",
                "bullets": bullets,
                "media": {"path": media.name},
            }
        ],
    )
    assert result["slide_count"] > 1
    rendered = "\n".join(slide_texts(presentation))
    assert [rendered.count(f"SENTINEL-{index:02d}") for index in range(12)] == [1] * 12
    assert rendered.count("ONE-CAPTION") == 1
    for slide in presentation.slides:
        assert len([shape for shape in slide.shapes if shape.name == "Penny Media"]) == 1
        assert len([shape for shape in slide.shapes if shape.name == "Penny Mark"]) == 1
    records = result["resolved_design"]["slides"]
    assert [record["continuation_index"] for record in records] == list(range(len(records)))
    assert len({record["palette"]["surface"] for record in records}) == 1
    assert all(record["source_slide_index"] == 1 for record in records)


def test_markdown_inherits_deck_design_without_changing_classification(
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    result, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        design={"palette": {"canvas": "080B12", "surface": "111827"}},
        markdown="# Deck\n\nSubtitle\n\n## Body\n\nParagraph.",
    )
    assert result["layouts_used"] == {"title": 1, "content": 1}
    assert all(record["active"] for record in result["resolved_design"]["slides"])
    assert all(_shape_named(slide, "Penny Canvas") for slide in presentation.slides)


def test_legacy_generation_does_not_add_design_shapes_or_change_legacy_palette(
    tmp_path: Path, powerpointgen: ModuleType
) -> None:
    result, presentation = generate_deck(
        powerpointgen,
        tmp_path,
        slides=[{"layout": "content", "title": "Legacy", "body": "Body"}],
    )
    assert not any(shape.name.startswith("Penny ") for shape in presentation.slides[0].shapes)
    assert result["resolved_palette"]["background"] == "FFFFFF"
    assert result["resolved_design"]["slides"][0]["active"] is False
