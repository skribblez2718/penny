#!/usr/bin/env python3
"""Generate a professionally styled PowerPoint (.pptx) presentation.

Reads a JSON spec from stdin (structured `slides` list or `markdown`) and
renders a 16:9 deck through python-pptx.
"""

from __future__ import annotations

import io
import json
import os
import re
import sys
import tempfile
import unicodedata
import zipfile
from dataclasses import asdict, dataclass
from functools import lru_cache
from pathlib import Path
from typing import Any, Callable, Sequence

from lxml import etree
from markdown_it import MarkdownIt
from markdown_it.token import Token
from PIL import Image, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.presentation import Presentation as PresentationType
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ============================================================
# Theme system
# ============================================================


@dataclass(frozen=True)
class Theme:
    accent: str
    accent_light: str
    text_dark: str
    text_muted: str
    heading_font: str
    body_font: str
    mono_font: str


@dataclass(frozen=True)
class Palette:
    text: str
    text_muted: str
    heading: str
    accent: str
    accent_text: str
    accent_soft: str
    on_accent: str
    link: str
    link_on_accent: str
    background: str
    surface: str
    border: str
    code_background: str


@dataclass(frozen=True)
class FontChoice:
    role: str
    requested: str
    resolved: str
    substituted: bool
    verified: bool
    metrics_path: str | None
    metrics_face_index: int = 0
    metrics_styles: dict[str, tuple[str, int]] | None = None


@dataclass(frozen=True)
class FontMetricPlan:
    regular: str | tuple[str, int] | None
    styles: dict[str, tuple[str, int]]


FontCatalogEntry = tuple[str, str] | tuple[str, str, int]
FontMetricSource = str | tuple[str, int] | FontMetricPlan | None


@dataclass
class Options:
    theme_name: str
    theme: Theme
    palette: Palette
    font_plan: list[FontChoice]
    title: str | None
    subtitle: str | None
    author: str | None
    date: str | None
    footer_text: str | None
    slide_numbers: bool
    line_break_mode: str
    output_path: str
    staging_path: str | None
    project_root: str


THEMES: dict[str, Theme] = {
    "executive": Theme(
        "1F3A5F", "D9E2F3", "1F2937", "6B7280", "Calibri Light", "Calibri", "Consolas"
    ),
    "modern": Theme("4F46E5", "E0E7FF", "111827", "6B7280", "Segoe UI", "Segoe UI", "Consolas"),
    "minimal": Theme("111827", "E5E7EB", "111827", "6B7280", "Arial", "Arial", "Consolas"),
    "editorial": Theme("7C2D12", "EFDFD3", "1F2937", "6B7280", "Georgia", "Georgia", "Consolas"),
    "tech": Theme("0F766E", "CCFBF1", "111827", "6B7280", "Segoe UI", "Calibri", "Consolas"),
}

BAND_FILL = "F5F7FA"

LAYOUTS = ["title", "section", "content", "two_column", "table", "quote", "image", "closing"]

# Slide geometry (inches, 16:9)
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.7
CONTENT_W = SLIDE_W - 2 * MARGIN
CONTENT_TOP = 1.9
CONTENT_BOTTOM = 6.85

CONTENT_HEIGHT = CONTENT_BOTTOM - CONTENT_TOP
BODY_FONT_PT = 14.0
BODY_LINE_HEIGHT_IN = 0.25
PARAGRAPH_GAP_IN = 0.11
CODE_FONT_PT = 12.0
CODE_LINE_HEIGHT_IN = 0.24
CODE_PANEL_PADDING_IN = 0.25
TABLE_HEADER_FONT_PT = 13.0
TABLE_BODY_FONT_PT = 12.0
TABLE_LINE_HEIGHT_IN = 0.23
TABLE_CELL_VERTICAL_PADDING_IN = 0.12
TABLE_MIN_ROW_HEIGHT_IN = 0.38
TABLE_HEADER_MIN_HEIGHT_IN = 0.5
BULLET_FONT_TIERS = (
    (16.0, 14.0, 12.5),
    (14.0, 13.0, 12.0),
    (13.0, 12.0, 11.0),
)
BULLET_LINE_HEIGHT_FACTOR = 1.25
BULLET_GAP_FACTOR = 0.45
MIN_EFFECTIVE_PPI = 96.0
POINTS_PER_INCH = 72.0
PIXELS_PER_INCH = 96.0

_MD_INLINE = MarkdownIt("commonmark").enable(["strikethrough"])

_SANS_FALLBACKS = ["Arial", "Liberation Sans", "Noto Sans", "DejaVu Sans"]
_SERIF_FALLBACKS = ["Georgia", "Liberation Serif", "Noto Serif", "DejaVu Serif"]
_MONO_FALLBACKS = ["Consolas", "Liberation Mono", "Noto Sans Mono", "DejaVu Sans Mono"]
_SERIF_FAMILIES = {"georgia", "times", "times new roman"}
_FONT_SUFFIXES = {".ttf", ".otf", ".ttc"}

_HEX_RE = re.compile(r"^[0-9A-Fa-f]{6}$")
_BR_TAG_RE = re.compile(r"<br\s*/?>", re.IGNORECASE)
_HTML_IMAGE_RE = re.compile(r"<\s*img\b", re.IGNORECASE)
_PARAGRAPH_SPLIT_RE = re.compile(r"\n[ \t]*\n+")
_REQUIRED_PPTX_PARTS = {
    "[Content_Types].xml",
    "_rels/.rels",
    "docProps/core.xml",
    "docProps/app.xml",
    "ppt/presentation.xml",
    "ppt/_rels/presentation.xml.rels",
}


# ============================================================
# Color + contrast helpers
# ============================================================


def _mix(hex_color: str, other: str, factor: float) -> str:
    """Blend ``hex_color`` toward ``other`` by ``factor`` (0..1)."""

    a = [int(hex_color[i : i + 2], 16) for i in (0, 2, 4)]
    b = [int(other[i : i + 2], 16) for i in (0, 2, 4)]
    return "".join(f"{round(x + (y - x) * factor):02X}" for x, y in zip(a, b, strict=True))


def _relative_luminance(hex_color: str) -> float:
    channels = []
    for value in _hex_to_triplet(hex_color):
        normalized = value / 255.0
        if normalized <= 0.03928:
            channels.append(normalized / 12.92)
        else:
            channels.append(((normalized + 0.055) / 1.055) ** 2.4)
    return 0.2126 * channels[0] + 0.7152 * channels[1] + 0.0722 * channels[2]


def _hex_to_triplet(hex_color: str) -> tuple[int, int, int]:
    return (int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def _contrast_ratio(foreground: str, background: str) -> float:
    light = max(_relative_luminance(foreground), _relative_luminance(background))
    dark = min(_relative_luminance(foreground), _relative_luminance(background))
    return (light + 0.05) / (dark + 0.05)


def _ensure_contrast(foreground: str, background: str, minimum: float = 4.5) -> str:
    if _contrast_ratio(foreground, background) >= minimum:
        return foreground
    direction = (
        "000000"
        if _contrast_ratio("000000", background) > _contrast_ratio("FFFFFF", background)
        else "FFFFFF"
    )
    for step in range(1, 18):
        candidate = _mix(foreground, direction, step / 18.0)
        if _contrast_ratio(candidate, background) >= minimum:
            return candidate
    return direction


def _derive_palette(theme: Theme) -> Palette:
    accent = theme.accent
    on_accent = (
        "000000"
        if _contrast_ratio("000000", accent) >= _contrast_ratio("FFFFFF", accent)
        else "FFFFFF"
    )
    return Palette(
        text=_ensure_contrast(theme.text_dark, "FFFFFF"),
        text_muted=_ensure_contrast(theme.text_muted, "FFFFFF"),
        heading=_ensure_contrast(theme.text_dark, "FFFFFF"),
        accent=accent,
        accent_text=_ensure_contrast(accent, "FFFFFF"),
        accent_soft=_mix(accent, "FFFFFF", 0.88),
        on_accent=on_accent,
        link=_ensure_contrast(accent, "FFFFFF"),
        link_on_accent=on_accent,
        background="FFFFFF",
        surface="FFFFFF",
        border=_mix(theme.text_dark, "FFFFFF", 0.85),
        code_background=_mix(theme.text_dark, "FFFFFF", 0.95),
    )


# ============================================================
# Options
# ============================================================


def _opt_str(spec: dict[str, Any], key: str) -> str | None:
    value = spec.get(key)
    if value is None or value == "":
        return None
    return str(value)


def _opt_enum(spec: dict[str, Any], key: str, allowed: list[str], default: str) -> str:
    value = str(spec.get(key) or default).lower()
    if value not in allowed:
        raise ValueError(f"{key} must be one of {allowed}, got {value!r}")
    return value


def _resolve_theme(spec: dict[str, Any]) -> tuple[str, Theme]:
    name = str(spec.get("theme") or "executive").lower()
    if name not in THEMES:
        raise ValueError(f"theme must be one of {list(THEMES)}, got {name!r}")
    theme = THEMES[name]
    accent = _opt_str(spec, "accent_color")
    if accent:
        accent = accent.lstrip("#").upper()
        if not _HEX_RE.fullmatch(accent):
            raise ValueError(f"accent_color must be a 6-digit hex color, got {accent!r}")
        theme = Theme(
            accent,
            _mix(accent, "FFFFFF", 0.88),
            theme.text_dark,
            theme.text_muted,
            theme.heading_font,
            theme.body_font,
            theme.mono_font,
        )
    return name, theme


def _font_directories() -> list[Path]:
    directories = [
        Path("/usr/share/fonts"),
        Path("/usr/local/share/fonts"),
        Path.home() / ".fonts",
        Path.home() / ".local/share/fonts",
        Path("/System/Library/Fonts"),
        Path("/Library/Fonts"),
        Path.home() / "Library/Fonts",
    ]
    windows_root = os.environ.get("WINDIR")
    if windows_root:
        directories.append(Path(windows_root) / "Fonts")
    local_app_data = os.environ.get("LOCALAPPDATA")
    if local_app_data:
        directories.append(Path(local_app_data) / "Microsoft" / "Windows" / "Fonts")
    else:
        directories.append(Path.home() / "AppData" / "Local" / "Microsoft" / "Windows" / "Fonts")
    return sorted({directory for directory in directories if directory.is_dir()})


def _font_faces(path: Path) -> list[tuple[str, str, int]]:
    """Return every readable family/style face in a font file deterministically."""

    faces: list[tuple[str, str, int]] = []
    face_limit = 64 if path.suffix.lower() == ".ttc" else 1
    for index in range(face_limit):
        try:
            family, style = ImageFont.truetype(str(path), 12, index=index).getname()
        except (OSError, ValueError):
            break
        family_name = str(family).strip()
        if family_name:
            faces.append((family_name, str(style).strip(), index))
    return faces


def _font_style_rank(style: str) -> int:
    normalized = style.casefold().replace("-", " ")
    if normalized in {"regular", "normal", "roman"}:
        return 0
    if normalized in {"book", "medium"}:
        return 1
    return 2


def _font_style_name(style: str) -> str:
    normalized = style.casefold().replace("-", " ")
    is_bold = any(marker in normalized for marker in ("bold", "semibold", "demibold"))
    is_italic = any(marker in normalized for marker in ("italic", "oblique"))
    if is_bold and is_italic:
        return "bold_italic"
    if is_bold:
        return "bold"
    if is_italic:
        return "italic"
    return "regular"


def _font_style_choice_rank(style: str, style_name: str) -> int:
    normalized = style.casefold().replace("-", " ")
    if style_name == "regular":
        return _font_style_rank(style)
    exact = {
        "bold": {"bold"},
        "italic": {"italic", "oblique"},
        "bold_italic": {"bold italic", "bold oblique"},
    }
    return 0 if normalized in exact.get(style_name, set()) else 1


@lru_cache(maxsize=1)
def _font_style_catalog() -> dict[str, dict[str, tuple[str, int]]]:
    candidates: dict[str, dict[str, tuple[tuple[int, str, int], tuple[str, int]]]] = {}
    for directory in _font_directories():
        for path in sorted(directory.rglob("*")):
            if path.suffix.lower() not in _FONT_SUFFIXES:
                continue
            for family, style, face_index in _font_faces(path):
                family_key = family.casefold()
                style_name = _font_style_name(style)
                rank = (
                    _font_style_choice_rank(style, style_name),
                    str(path).casefold(),
                    face_index,
                )
                by_style = candidates.setdefault(family_key, {})
                current = by_style.get(style_name)
                source = (str(path), face_index)
                if current is None or rank < current[0]:
                    by_style[style_name] = (rank, source)
    return {
        family: {style: candidate[1] for style, candidate in by_style.items()}
        for family, by_style in candidates.items()
    }


@lru_cache(maxsize=1)
def _font_catalog() -> dict[str, FontCatalogEntry]:
    candidates: dict[str, tuple[tuple[int, str, int], FontCatalogEntry]] = {}
    for directory in _font_directories():
        for path in sorted(directory.rglob("*")):
            if path.suffix.lower() not in _FONT_SUFFIXES:
                continue
            for family, style, face_index in _font_faces(path):
                key = family.casefold()
                rank = (_font_style_rank(style), str(path).casefold(), face_index)
                entry: FontCatalogEntry = (family, str(path), face_index)
                current = candidates.get(key)
                if current is None or rank < current[0]:
                    candidates[key] = (rank, entry)
    return {family: candidate[1] for family, candidate in candidates.items()}


def _dedupe(values: list[str]) -> list[str]:
    seen: set[str] = set()
    result: list[str] = []
    for value in values:
        key = value.casefold()
        if key in seen:
            continue
        seen.add(key)
        result.append(value)
    return result


def _font_fallbacks(role: str, requested: str) -> list[str]:
    if role == "mono":
        return _MONO_FALLBACKS
    if requested.casefold() in _SERIF_FAMILIES:
        return _SERIF_FALLBACKS
    return _SANS_FALLBACKS


def _resolve_font(
    role: str,
    requested: str,
    catalog: dict[str, FontCatalogEntry] | None = None,
) -> FontChoice:
    available = _font_catalog() if catalog is None else catalog
    styles = _font_style_catalog() if catalog is None else {}
    candidates = _dedupe([requested, *_font_fallbacks(role, requested)])
    for candidate in candidates:
        match = available.get(candidate.casefold())
        if match is not None:
            resolved, metrics_path = match[:2]
            metrics_face_index = match[2] if len(match) == 3 else 0
            metrics_styles = dict(styles.get(resolved.casefold(), {}))
            metrics_styles.setdefault("regular", (metrics_path, metrics_face_index))
            if catalog is None and not {
                "regular",
                "bold",
                "italic",
                "bold_italic",
            } <= set(metrics_styles):
                continue
            return FontChoice(
                role=role,
                requested=requested,
                resolved=resolved,
                substituted=resolved.casefold() != requested.casefold(),
                verified=True,
                metrics_path=metrics_path,
                metrics_face_index=metrics_face_index,
                metrics_styles=metrics_styles,
            )
    return FontChoice(
        role=role,
        requested=requested,
        resolved=requested,
        substituted=False,
        verified=False,
        metrics_path=None,
    )


def parse_options(spec: dict[str, Any]) -> Options:
    theme_name, theme = _resolve_theme(spec)
    output_path = _opt_str(spec, "output_path")
    if not output_path:
        raise ValueError("output_path is required in the generator spec")

    heading_choice = _resolve_font("heading", theme.heading_font)
    body_choice = _resolve_font("body", theme.body_font)
    mono_choice = _resolve_font("mono", theme.mono_font)
    resolved_theme = Theme(
        accent=theme.accent,
        accent_light=theme.accent_light,
        text_dark=theme.text_dark,
        text_muted=theme.text_muted,
        heading_font=heading_choice.resolved,
        body_font=body_choice.resolved,
        mono_font=mono_choice.resolved,
    )

    staging = _opt_str(spec, "staging_path")
    return Options(
        theme_name=theme_name,
        theme=resolved_theme,
        palette=_derive_palette(resolved_theme),
        font_plan=[heading_choice, body_choice, mono_choice],
        title=_opt_str(spec, "title"),
        subtitle=_opt_str(spec, "subtitle"),
        author=_opt_str(spec, "author"),
        date=_opt_str(spec, "date"),
        footer_text=_opt_str(spec, "footer_text"),
        slide_numbers=bool(spec.get("slide_numbers", True)),
        line_break_mode=_opt_enum(spec, "line_break_mode", ["preserve", "commonmark"], "preserve"),
        output_path=output_path,
        staging_path=staging,
        project_root=_opt_str(spec, "project_root") or os.getcwd(),
    )


# ============================================================
# Slide normalization
# ============================================================


def _normalize_bullet(item: Any) -> dict[str, Any]:
    if isinstance(item, str):
        return {"text": item, "level": 0, "bold": False}
    if isinstance(item, dict) and "text" in item:
        level = int(item.get("level") or 0)
        return {
            "text": str(item["text"]),
            "level": max(0, min(level, 2)),
            "bold": bool(item.get("bold")),
        }
    raise ValueError(f"invalid bullet item: {item!r}")


def _body_parts(value: Any) -> list[str]:
    if value is None:
        return []
    return [part for part in _PARAGRAPH_SPLIT_RE.split(str(value)) if part]


def _normalize_column(column: Any) -> dict[str, Any]:
    if not isinstance(column, dict):
        return {"heading": None, "body": None, "body_parts": [], "bullets": []}
    return {
        "heading": column.get("heading"),
        "body": column.get("body"),
        "body_parts": _body_parts(column.get("body")),
        "bullets": [_normalize_bullet(b) for b in column.get("bullets") or []],
    }


def normalize_slide(raw: Any) -> dict[str, Any]:
    if not isinstance(raw, dict):
        raise ValueError(f"slide must be an object, got {type(raw).__name__}")
    layout = str(raw.get("layout") or "content")
    if layout not in LAYOUTS:
        raise ValueError(f"layout must be one of {LAYOUTS}, got {layout!r}")
    slide = dict(raw)
    slide["layout"] = layout
    if "body_parts" not in slide and slide.get("body") is not None:
        slide["body_parts"] = _body_parts(slide.get("body"))
    slide["body"] = str(slide["body"]) if slide.get("body") is not None else None
    slide["bullets"] = [_normalize_bullet(b) for b in raw.get("bullets") or []]
    if layout == "two_column":
        slide["left"] = _normalize_column(raw.get("left"))
        slide["right"] = _normalize_column(raw.get("right"))
    return slide


# ============================================================
# Markdown normalization
# ============================================================


def _validate_inline_html(content: str) -> bool:
    """Return True for the one safe supported HTML form; reject everything else."""

    if _BR_TAG_RE.fullmatch(content.strip()):
        return True
    if _HTML_IMAGE_RE.search(content):
        raise ValueError(
            "inline image in raw HTML is not supported; place the image in its own paragraph"
        )
    raise ValueError("unsupported inline HTML is not supported on slides")


def _inline_markup(inline: Token, context: str) -> str:
    children = inline.children or []
    if _contains_image(children):
        raise ValueError(
            f"inline image mixed with text is not supported in {context}; "
            "place the image in its own paragraph"
        )
    for child in children:
        if child.type == "html_inline":
            _validate_inline_html(child.content)
    return inline.content


def _child_href(token: Token) -> str:
    attrs = token.attrs or {}
    if isinstance(attrs, dict):
        return str(attrs.get("href") or "")
    for key, value in attrs:
        if key == "href":
            return str(value)
    return ""


def _is_image_only(children: list[Token]) -> tuple[bool, Token | None]:
    meaningful = [
        token for token in children if not (token.type == "text" and token.content.strip() == "")
    ]
    if len(meaningful) == 1 and meaningful[0].type == "image":
        return True, meaningful[0]
    return False, None


def _contains_image(children: list[Token]) -> bool:
    return any(token.type == "image" for token in children)


class _MarkdownSlicer:
    """Walks block tokens and accumulates slide dicts per documented rules."""

    def __init__(self, meta: Options) -> None:
        self.meta = meta
        self.slides: list[dict[str, Any]] = []
        self.current: dict[str, Any] | None = None
        self.section_pending = False
        self.warnings: list[str] = []

    @staticmethod
    def _blank(title: str | None = None) -> dict[str, Any]:
        return {
            "title": title,
            "bullets": [],
            "body_parts": [],
            "images": [],
            "code_parts": [],
            "tables": [],
        }

    def slice(self, tokens: list[Token]) -> list[dict[str, Any]]:
        i = self._maybe_title_slide(tokens)
        while i < len(tokens):
            i = self._consume(tokens, i)
        self._finalize()
        if not self.slides or self.slides[0]["layout"] != "title":
            if self.meta.title:
                self.slides.insert(0, {"layout": "title", "title": self.meta.title})
        return self.slides

    def _maybe_title_slide(self, tokens: list[Token]) -> int:
        if tokens and tokens[0].type == "heading_open" and tokens[0].tag == "h1":
            slide: dict[str, Any] = {
                "layout": "title",
                "title": _inline_markup(tokens[1], "title"),
            }
            i = 3
            if i < len(tokens) and tokens[i].type == "paragraph_open":
                children = tokens[i + 1].children or []
                only_image = _is_image_only(children)[0]
                if not only_image:
                    subtitle = _inline_markup(tokens[i + 1], "subtitle")
                    if subtitle:
                        slide["subtitle"] = subtitle
                        i += 3
            self.slides.append(slide)
            return i
        return 0

    def _consume(self, tokens: list[Token], i: int) -> int:
        token = tokens[i]
        if token.type == "hr":
            self.section_pending = True
            return i + 1
        if token.type == "heading_open":
            return self._heading(tokens, i)
        self.section_pending = False
        if token.type == "paragraph_open":
            return self._paragraph(tokens, i)
        if token.type in ("bullet_list_open", "ordered_list_open"):
            return self._list(tokens, i, 0)
        if token.type == "blockquote_open":
            return self._blockquote(tokens, i)
        if token.type == "table_open":
            return self._table(tokens, i)
        if token.type in ("fence", "code_block"):
            self._slide()["code_parts"].append(token.content)
            return i + 1
        if token.type == "html_block":
            if _HTML_IMAGE_RE.search(token.content):
                raise ValueError(
                    "inline image in raw HTML is not supported; "
                    "place the image in its own paragraph"
                )
            raise ValueError("unsupported HTML block is not supported on slides")
        return i + 1

    def _slide(self) -> dict[str, Any]:
        if self.current is None:
            self.current = self._blank()
        return self.current

    def _heading(self, tokens: list[Token], i: int) -> int:
        level = int(tokens[i].tag[1])
        text = _inline_markup(tokens[i + 1], "heading")
        # Any heading consumes a pending hr, so only an H2 IMMEDIATELY after the
        # hr becomes a section divider.
        pending, self.section_pending = self.section_pending, False
        if level <= 2:
            self._finalize()
            if pending and level == 2:
                self.slides.append({"layout": "section", "title": text})
            else:
                self.current = self._blank(text)
        else:
            self._slide()["bullets"].append({"text": text, "level": 0, "bold": True})
        return i + 3

    def _paragraph(self, tokens: list[Token], i: int) -> int:
        children = tokens[i + 1].children or []
        only_image, image_token = _is_image_only(children)
        if only_image:
            assert image_token is not None
            self._slide()["images"].append(
                {
                    "path": str(image_token.attrs.get("src", "")),
                    "caption": image_token.content or None,
                }
            )
        else:
            if _contains_image(children):
                raise ValueError(
                    "inline image mixed with surrounding text is not supported on slides"
                )
            self._slide()["body_parts"].append(tokens[i + 1].content)
        return i + 3

    def _list(self, tokens: list[Token], i: int, depth: int) -> int:
        close = tokens[i].type.replace("_open", "_close")
        i += 1
        while tokens[i].type != close:
            if tokens[i].type == "list_item_open":
                i = self._list_item(tokens, i, depth)
            else:
                i += 1
        return i + 1

    def _list_item(self, tokens: list[Token], i: int, depth: int) -> int:
        i += 1
        while tokens[i].type != "list_item_close":
            if tokens[i].type == "paragraph_open":
                children = tokens[i + 1].children or []
                only_image, image_token = _is_image_only(children)
                if only_image:
                    raise ValueError(
                        "image-only bullet items are not supported in markdown list content"
                    )
                if _contains_image(children):
                    raise ValueError(
                        "inline image mixed with surrounding text is not supported on slides"
                    )
                self._slide()["bullets"].append(
                    {"text": tokens[i + 1].content, "level": min(depth, 2), "bold": False}
                )
                i += 3
            elif tokens[i].type in ("bullet_list_open", "ordered_list_open"):
                i = self._list(tokens, i, depth + 1)
            else:
                i += 1
        return i + 1

    def _blockquote(self, tokens: list[Token], i: int) -> int:
        parts: list[str] = []
        while tokens[i].type != "blockquote_close":
            if tokens[i].type == "inline":
                parts.append(tokens[i].content)
            i += 1
        text = "\n".join(parts).strip()
        attribution = None
        lines = text.split("\n")
        if len(lines) > 1 and lines[-1].lstrip().startswith(("—", "--")):
            attribution = lines[-1].lstrip("—-").strip()
            text = "\n".join(lines[:-1]).strip()
        slide = self._slide()
        if "quote" not in slide and not slide["bullets"] and not slide["body_parts"]:
            slide["quote"] = text
            slide["attribution"] = attribution
        else:
            # A second quote (or a quote after other content) joins the body instead of overwriting.
            self._slide()["body_parts"].append(text)
        return i + 1

    def _table(self, tokens: list[Token], i: int) -> int:
        headers: list[str] = []
        rows: list[list[str]] = []
        row: list[str] = []
        in_head = False
        while tokens[i].type != "table_close":
            if tokens[i].type == "thead_open":
                in_head = True
            elif tokens[i].type == "thead_close":
                in_head = False
            elif tokens[i].type == "tr_open":
                row = []
            elif tokens[i].type == "inline":
                row.append(tokens[i].content)
            elif tokens[i].type == "tr_close":
                if in_head:
                    headers = row
                else:
                    rows.append(row)
            i += 1
        self._slide()["tables"].append({"headers": headers, "rows": rows})
        return i + 1

    def _finalize(self) -> None:
        slide = self.current
        self.current = None
        if slide is None:
            return
        self.slides.extend(self._materialize(slide))

    def _materialize(self, slide: dict[str, Any]) -> list[dict[str, Any]]:
        title = slide.get("title")
        body_parts = [str(part) for part in (slide.get("body_parts") or [])]
        bullets = slide.get("bullets") or []
        tables = list(slide.get("tables") or [])
        code_parts = slide.get("code_parts") or []

        pages: list[dict[str, Any]] = []
        if slide.get("quote"):
            pages.append(
                {
                    "layout": "quote",
                    "title": title,
                    "quote": slide["quote"],
                    "attribution": slide.get("attribution"),
                }
            )

        body_on_table = bool(tables) and not bullets and not code_parts
        if bullets or code_parts or (body_parts and not body_on_table):
            pages.append(
                {
                    "layout": "content",
                    "title": title if not pages else None,
                    "kicker": None,
                    "body_parts": body_parts,
                    "body": "\n\n".join(body_parts) if body_parts else None,
                    "code_parts": [str(code) for code in code_parts],
                    "code": [str(code) for code in code_parts],
                    "bullets": bullets,
                    "code_lines": None,
                }
            )

        for table_index, table in enumerate(tables):
            table_body_parts = body_parts if body_on_table and table_index == 0 else []
            pages.append(
                {
                    "layout": "table",
                    "title": title if not pages else None,
                    "table": table,
                    "body_parts": table_body_parts,
                    "body": "\n\n".join(table_body_parts) if table_body_parts else None,
                }
            )

        for image in slide.get("images") or []:
            pages.append(
                {
                    "layout": "image",
                    "title": title if not pages else None,
                    "image_path": image["path"],
                    "caption": image["caption"],
                }
            )

        if not pages and title:
            pages.append({"layout": "content", "title": title, "body_parts": [], "bullets": []})
        return pages


def markdown_to_slides(markdown: str, meta: Options) -> tuple[list[dict[str, Any]], list[str]]:
    parser = MarkdownIt("commonmark").enable(["table", "strikethrough"])
    tokens = parser.parse(markdown)
    slicer = _MarkdownSlicer(meta)
    slides = [normalize_slide(s) for s in slicer.slice(tokens)]
    return slides, slicer.warnings


# ============================================================
# Pagination helpers
# ============================================================


def _visible_lines(text: str, line_break_mode: str) -> list[str]:
    tokens = _MD_INLINE.parseInline(text)
    children = tokens[0].children if tokens else []
    segments: list[str] = []
    current = ""
    for token in children or []:
        if token.type in ("text", "code_inline"):
            current += token.content
        elif token.type == "softbreak":
            if line_break_mode == "preserve":
                segments.append(current)
                current = ""
            else:
                current += " "
        elif token.type == "hardbreak":
            segments.append(current)
            current = ""
        elif token.type == "html_inline":
            if _validate_inline_html(token.content):
                segments.append(current)
                current = ""
        elif token.type == "image":
            raise ValueError(
                "inline image mixed with text is not supported; "
                "place the image in its own paragraph"
            )
    segments.append(current)
    return segments or [""]


def _fallback_visual_units(text: str) -> float:
    units = 0.0
    for character in text:
        if unicodedata.combining(character):
            continue
        units += 2.0 if unicodedata.east_asian_width(character) in {"W", "F"} else 1.0
    return units


def _font_for_metrics(
    source: FontMetricSource,
    font_pt: float,
    style: str = "regular",
) -> ImageFont.FreeTypeFont | None:
    if not source:
        return None
    selected: str | tuple[str, int] | None
    if isinstance(source, FontMetricPlan):
        selected = source.regular if style == "regular" else source.styles.get(style)
        if selected is None:
            return None
    else:
        selected = source
    font_path, face_index = (selected, 0) if isinstance(selected, str) else selected
    try:
        pixel_size = max(1, round(font_pt * PIXELS_PER_INCH / POINTS_PER_INCH))
        return ImageFont.truetype(font_path, pixel_size, index=face_index)
    except (OSError, ValueError):
        return None


def _wrap_overwide_token(
    token: str,
    current_width: float,
    max_width: float,
    measure: Callable[[str], float],
) -> tuple[int, float]:
    completed_lines = 1 if current_width else 0
    current_width = 0.0
    for character in token:
        character_width = float(measure(character))
        if current_width and current_width + character_width > max_width:
            completed_lines += 1
            current_width = 0.0
        if character_width > max_width:
            completed_lines += 1
        else:
            current_width += character_width
    return completed_lines, current_width


def _segment_line_count(
    segment: str,
    width_in: float,
    font_pt: float,
    font_source: FontMetricSource,
) -> int:
    """Greedily simulate word wrapping, with character fallback for wide tokens."""

    if not segment:
        return 1
    font = _font_for_metrics(font_source, font_pt)
    max_width = max(1.0, width_in * PIXELS_PER_INCH)
    average_character_width = font_pt * 0.58 * PIXELS_PER_INCH / POINTS_PER_INCH

    def measure(value: str) -> float:
        if font is not None:
            return float(font.getlength(value))
        return _fallback_visual_units(value) * average_character_width

    line_count = 0
    current_width = 0.0
    pending_space = 0.0
    for token in re.findall(r"\S+|[ \t]+", segment):
        if token.isspace():
            pending_space = measure(token)
            continue

        word_width = measure(token)
        separator_width = pending_space if current_width else 0.0
        pending_space = 0.0
        if word_width <= max_width:
            if current_width and current_width + separator_width + word_width > max_width:
                line_count += 1
                current_width = word_width
            else:
                current_width += separator_width + word_width
            continue

        completed, current_width = _wrap_overwide_token(
            token,
            current_width,
            max_width,
            measure,
        )
        line_count += completed

    return max(1, line_count + (1 if current_width else 0))


def _wrapped_measured_line_count(
    tokens: list[tuple[float, float, list[float]]],
    max_width: float,
) -> int:
    line_count = 0
    current_width = 0.0
    for space_width, word_width, character_widths in tokens:
        separator = space_width if current_width else 0.0
        if word_width <= max_width:
            if current_width and current_width + separator + word_width > max_width:
                line_count += 1
                current_width = word_width
            else:
                current_width += separator + word_width
            continue

        if current_width:
            line_count += 1
            current_width = 0.0
        for character_width in character_widths:
            if current_width and current_width + character_width > max_width:
                line_count += 1
                current_width = 0.0
            if character_width > max_width:
                line_count += 1
            else:
                current_width += character_width
    return max(1, line_count + (1 if current_width else 0))


def _estimate_lines(  # noqa: C901
    text: str,
    width_in: float,
    font_pt: float,
    line_break_mode: str,
    font_path: FontMetricSource = None,
    mono_font_path: FontMetricSource = None,
    bold: bool = False,
    italic: bool = False,
) -> int:
    """Measure Markdown runs with the same context/mono face split used by rendering."""

    body_source = font_path
    mono_source = mono_font_path or font_path
    font_cache: dict[tuple[bool, str], ImageFont.FreeTypeFont | None] = {}
    average_width = font_pt * 0.58 * PIXELS_PER_INCH / POINTS_PER_INCH
    mono_average_width = font_pt * 0.62 * PIXELS_PER_INCH / POINTS_PER_INCH
    logical_lines: list[list[tuple[float, float, list[float]]]] = []
    line_tokens: list[tuple[float, float, list[float]]] = []
    word_widths: list[float] = []
    pending_space_width = 0.0
    state = {"bold": bold, "italic": italic}

    def measure(character: str, code: bool) -> float:
        if state["bold"] and state["italic"]:
            style = "bold_italic"
        elif state["bold"]:
            style = "bold"
        elif state["italic"]:
            style = "italic"
        else:
            style = "regular"
        source = mono_source if code else body_source
        cache_key = (code, style)
        if cache_key not in font_cache:
            font_cache[cache_key] = _font_for_metrics(source, font_pt, style)
        font = font_cache[cache_key]
        if font is None and style != "regular" and isinstance(source, FontMetricPlan):
            raise ValueError(
                f"verified {style.replace('_', '-')} font metrics are unavailable; "
                "cannot preflight styled text safely"
            )
        if font is not None:
            return float(font.getlength(character))
        return _fallback_visual_units(character) * (mono_average_width if code else average_width)

    def flush_word() -> None:
        nonlocal pending_space_width
        if word_widths:
            line_tokens.append((pending_space_width, sum(word_widths), list(word_widths)))
            word_widths.clear()
            pending_space_width = 0.0

    def line_break() -> None:
        nonlocal line_tokens, pending_space_width
        flush_word()
        logical_lines.append(line_tokens)
        line_tokens = []
        pending_space_width = 0.0

    tokens = _MD_INLINE.parseInline(text)
    children = tokens[0].children if tokens else []
    for token in children or []:
        if token.type in {"strong_open", "strong_close"}:
            state["bold"] = True if token.type == "strong_open" else bold
            continue
        if token.type in {"em_open", "em_close"}:
            state["italic"] = True if token.type == "em_open" else italic
            continue
        if token.type == "softbreak":
            if line_break_mode == "preserve":
                line_break()
            else:
                flush_word()
                pending_space_width += measure(" ", False)
            continue
        if token.type == "hardbreak":
            line_break()
            continue
        if token.type == "html_inline":
            if _validate_inline_html(token.content):
                line_break()
            continue
        if token.type == "image":
            raise ValueError(
                "inline image mixed with text is not supported; "
                "place the image in its own paragraph"
            )
        if token.type not in {"text", "code_inline"}:
            continue
        is_code = token.type == "code_inline"
        for character in token.content:
            if character.isspace():
                flush_word()
                pending_space_width += measure(character, is_code)
            else:
                word_widths.append(measure(character, is_code))

    flush_word()
    logical_lines.append(line_tokens)
    max_width = max(1.0, width_in * PIXELS_PER_INCH)
    return max(
        1,
        sum(_wrapped_measured_line_count(line, max_width) for line in logical_lines),
    )


def _plain_line_count(
    text: str,
    width_in: float,
    font_pt: float,
    font_path: FontMetricSource,
) -> int:
    return max(
        1,
        sum(
            _segment_line_count(line, width_in, font_pt, font_path)
            for line in text.splitlines() or [""]
        ),
    )


def _body_paragraph_height(
    text: str,
    line_break_mode: str,
    font_path: FontMetricSource,
    mono_font_path: FontMetricSource = None,
    width_in: float = CONTENT_W,
    font_pt: float = BODY_FONT_PT,
    line_height_in: float = BODY_LINE_HEIGHT_IN,
) -> float:
    lines = _estimate_lines(
        text,
        width_in,
        font_pt,
        line_break_mode,
        font_path,
        mono_font_path,
    )
    return lines * line_height_in + PARAGRAPH_GAP_IN


def _bullet_height(
    bullet: dict[str, Any],
    sizes: tuple[float, float, float],
    line_break_mode: str,
    font_path: FontMetricSource,
    mono_font_path: FontMetricSource = None,
    width_in: float = CONTENT_W * 0.95,
) -> float:
    level = int(bullet.get("level") or 0)
    font_pt = sizes[level]
    usable_width = max(0.5, width_in - level * 0.28)
    lines = _estimate_lines(
        str(bullet.get("text", "")),
        usable_width,
        font_pt,
        line_break_mode,
        font_path,
        mono_font_path,
        bold=bool(bullet.get("bold")),
    )
    line_height = font_pt * BULLET_LINE_HEIGHT_FACTOR / POINTS_PER_INCH
    gap = max(4.0, font_pt * BULLET_GAP_FACTOR) / POINTS_PER_INCH
    return lines * line_height + gap


def _code_source_line_height(text: str, font_path: FontMetricSource) -> float:
    lines = _plain_line_count(text, CONTENT_W - 0.3, CODE_FONT_PT, font_path)
    return lines * CODE_LINE_HEIGHT_IN


def _table_row_height(
    row: Sequence[str],
    col_width_in: float,
    line_break_mode: str,
    font_path: FontMetricSource,
    mono_font_path: FontMetricSource = None,
    font_pt: float = TABLE_BODY_FONT_PT,
    minimum: float = TABLE_MIN_ROW_HEIGHT_IN,
    bold: bool = False,
) -> float:
    lines = max(
        1,
        max(
            (
                _estimate_lines(
                    str(value),
                    max(0.3, col_width_in - 0.24),
                    font_pt,
                    line_break_mode,
                    font_path,
                    mono_font_path,
                    bold=bold,
                )
                for value in row
            ),
            default=1,
        ),
    )
    return max(minimum, lines * TABLE_LINE_HEIGHT_IN + TABLE_CELL_VERTICAL_PADDING_IN)


def _paginate_content(  # noqa: C901
    title: str | None,
    kicker: str | None,
    body_parts: list[str],
    bullets: list[dict[str, Any]],
    code_parts: list[str],
    line_break_mode: str,
    body_font_path: FontMetricSource,
    mono_font_path: FontMetricSource,
) -> list[dict[str, Any]]:
    code_lines = [line.rstrip("\r") for block in code_parts for line in block.split("\n")]
    if code_lines and code_lines[-1] == "" and any(block.endswith("\n") for block in code_parts):
        code_lines.pop()

    if not body_parts and not bullets and not code_lines:
        return [
            {
                "layout": "content",
                "title": title,
                "kicker": kicker,
                "body_parts": [],
                "body": None,
                "code_lines": [],
                "code": [],
                "code_parts": [],
                "bullets": [],
                "_bullet_sizes": BULLET_FONT_TIERS[0],
            }
        ]

    pages: list[dict[str, Any]] = []
    body_cursor = 0
    code_cursor = 0
    bullet_cursor = 0

    while (
        body_cursor < len(body_parts)
        or code_cursor < len(code_lines)
        or bullet_cursor < len(bullets)
    ):
        used_height = 0.0
        page_body: list[str] = []
        page_code: list[str] = []
        page_bullets: list[dict[str, Any]] = []
        bullet_sizes = BULLET_FONT_TIERS[0]
        active_class_blocked = False

        while body_cursor < len(body_parts):
            paragraph = str(body_parts[body_cursor])
            paragraph_height = _body_paragraph_height(
                paragraph,
                line_break_mode,
                body_font_path,
                mono_font_path,
            )
            if paragraph_height > CONTENT_HEIGHT:
                raise ValueError("single paragraph exceeds content area")
            if used_height + paragraph_height > CONTENT_HEIGHT:
                active_class_blocked = True
                break
            page_body.append(paragraph)
            used_height += paragraph_height
            body_cursor += 1

        if not active_class_blocked:
            while code_cursor < len(code_lines):
                line = code_lines[code_cursor]
                line_height = _code_source_line_height(line, mono_font_path)
                panel_padding = CODE_PANEL_PADDING_IN if not page_code else 0.0
                if line_height + CODE_PANEL_PADDING_IN > CONTENT_HEIGHT:
                    raise ValueError("single code line exceeds content area")
                if used_height + panel_padding + line_height > CONTENT_HEIGHT:
                    active_class_blocked = True
                    break
                page_code.append(line)
                used_height += panel_padding + line_height
                code_cursor += 1

        if not active_class_blocked and bullet_cursor < len(bullets):
            first_bullet = bullets[bullet_cursor]
            fitting_tiers = [
                sizes
                for sizes in BULLET_FONT_TIERS
                if _bullet_height(
                    first_bullet,
                    sizes,
                    line_break_mode,
                    body_font_path,
                    mono_font_path,
                )
                <= CONTENT_HEIGHT - used_height
            ]
            if not fitting_tiers:
                fits_clean_page = any(
                    _bullet_height(
                        first_bullet,
                        sizes,
                        line_break_mode,
                        body_font_path,
                        mono_font_path,
                    )
                    <= CONTENT_HEIGHT
                    for sizes in BULLET_FONT_TIERS
                )
                if fits_clean_page and (page_body or page_code):
                    active_class_blocked = True
                else:
                    raise ValueError("single bullet item exceeds content area")
            else:
                bullet_sizes = fitting_tiers[0]

        if not active_class_blocked:
            while bullet_cursor < len(bullets):
                bullet = bullets[bullet_cursor]
                item_height = _bullet_height(
                    bullet,
                    bullet_sizes,
                    line_break_mode,
                    body_font_path,
                    mono_font_path,
                )
                if used_height + item_height > CONTENT_HEIGHT:
                    break
                page_bullets.append(bullet)
                used_height += item_height
                bullet_cursor += 1

        if not (page_body or page_code or page_bullets):
            raise ValueError("content cannot be placed without violating readability limits")

        page_index = len(pages)
        page_title = title if page_index == 0 else (f"{title} (cont.)" if title else None)
        pages.append(
            {
                "layout": "content",
                "title": page_title,
                "kicker": kicker,
                "body_parts": page_body,
                "body": "\n\n".join(page_body) if page_body else None,
                "code_lines": page_code,
                "code": ["\n".join(page_code)] if page_code else [],
                "code_parts": ["\n".join(page_code)] if page_code else [],
                "bullets": page_bullets,
                "_bullet_sizes": bullet_sizes,
            }
        )

    return pages


def _paginate_table(
    title: str | None,
    kicker: str | None,
    body_parts: list[str],
    table: dict[str, Any],
    line_break_mode: str,
    body_font_path: FontMetricSource,
    mono_font_path: FontMetricSource,
) -> list[dict[str, Any]]:
    headers = [str(header) for header in (table.get("headers") or [])]
    rows = [[str(cell) for cell in row] for row in (table.get("rows") or [])]
    column_count = max(1, len(headers), max((len(row) for row in rows), default=1))
    column_width = CONTENT_W / column_count
    header_height = _table_row_height(
        headers,
        column_width,
        line_break_mode,
        body_font_path,
        mono_font_path,
        font_pt=TABLE_HEADER_FONT_PT,
        minimum=TABLE_HEADER_MIN_HEIGHT_IN,
        bold=True,
    )
    if header_height > CONTENT_HEIGHT:
        raise ValueError("table header exceeds available table area")

    intro_parts = [str(part) for part in body_parts]
    intro_height = sum(
        _body_paragraph_height(
            part,
            line_break_mode,
            body_font_path,
            mono_font_path,
        )
        for part in intro_parts
    )
    first_row_height = (
        _table_row_height(
            rows[0],
            column_width,
            line_break_mode,
            body_font_path,
            mono_font_path,
        )
        if rows
        else 0.0
    )
    intro_pages: list[dict[str, Any]] = []
    if intro_parts and intro_height + header_height + first_row_height > CONTENT_HEIGHT:
        intro_pages = _paginate_content(
            title=title,
            kicker=kicker,
            body_parts=intro_parts,
            bullets=[],
            code_parts=[],
            line_break_mode=line_break_mode,
            body_font_path=body_font_path,
            mono_font_path=mono_font_path,
        )
        intro_parts = []
        intro_height = 0.0

    if not rows:
        global_index = len(intro_pages)
        page_title = title if global_index == 0 else (f"{title} (cont.)" if title else None)
        return [
            *intro_pages,
            {
                "layout": "table",
                "title": page_title,
                "kicker": kicker,
                "body_parts": intro_parts,
                "body": "\n\n".join(intro_parts) if intro_parts else None,
                "table": {
                    "headers": headers,
                    "rows": [],
                    "header_height": header_height,
                    "row_heights": [],
                },
            },
        ]

    pages: list[dict[str, Any]] = list(intro_pages)
    row_cursor = 0
    table_page_index = 0
    while row_cursor < len(rows):
        available_height = CONTENT_HEIGHT - header_height
        if table_page_index == 0:
            available_height -= intro_height
        chunk: list[list[str]] = []
        row_heights: list[float] = []
        used_height = 0.0

        while row_cursor < len(rows):
            row = rows[row_cursor]
            row_height = _table_row_height(
                row,
                column_width,
                line_break_mode,
                body_font_path,
                mono_font_path,
            )
            if row_height > CONTENT_HEIGHT - header_height:
                raise ValueError("single table row exceeds available table area")
            if used_height + row_height > available_height:
                break
            chunk.append(row)
            row_heights.append(row_height)
            used_height += row_height
            row_cursor += 1

        if not chunk:
            raise ValueError("single table row exceeds available table area")

        global_index = len(pages)
        page_title = title if global_index == 0 else (f"{title} (cont.)" if title else None)
        pages.append(
            {
                "layout": "table",
                "title": page_title,
                "kicker": kicker,
                "body_parts": intro_parts if table_page_index == 0 else [],
                "body": "\n\n".join(intro_parts) if table_page_index == 0 else None,
                "table": {
                    "headers": headers,
                    "rows": chunk,
                    "header_height": header_height,
                    "row_heights": row_heights,
                },
            }
        )
        table_page_index += 1

    return pages


# ============================================================
# Text helpers
# ============================================================


def _add_run(
    paragraph: Any,
    text: str,
    font: str,
    size_pt: float,
    color: str,
    state: dict[str, bool],
    href: str | None,
    link_color: str,
) -> None:
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.color.rgb = RGBColor.from_string(link_color if href else color)
    run.font.bold = state.get("bold", False)
    run.font.italic = state.get("italic", False)
    if state.get("strike"):
        run.font._rPr.set("strike", "sngStrike")
    if href:
        run.hyperlink.address = href
        run.font.underline = True


def _add_md_runs(  # noqa: C901
    paragraph: Any,
    text: str,
    font: str,
    size_pt: float,
    color: str,
    mono_font: str,
    line_break_mode: str,
    bold: bool = False,
    italic: bool = False,
    link_color: str | None = None,
) -> None:
    """Render inline markdown emphasis (**bold**, *italic*, `code`, ~~strike~~) as runs."""
    tokens = _MD_INLINE.parseInline(text)
    children = tokens[0].children if tokens else []
    state = {"bold": bold, "italic": italic, "strike": False}
    href: str | None = None
    toggles = {
        "strong_open": ("bold", True),
        "strong_close": ("bold", bold),
        "em_open": ("italic", True),
        "em_close": ("italic", italic),
        "s_open": ("strike", True),
        "s_close": ("strike", False),
    }

    resolved_link_color = link_color or color

    def add_linebreak() -> None:
        paragraph.add_line_break()

    for token in children or []:
        if token.type in toggles:
            key, value = toggles[token.type]
            state[key] = value
            continue
        if token.type == "link_open":
            href = _child_href(token)
            continue
        if token.type == "link_close":
            href = None
            continue
        if token.type == "softbreak":
            if line_break_mode == "preserve":
                add_linebreak()
            else:
                _add_run(
                    paragraph,
                    " ",
                    font,
                    size_pt,
                    color,
                    state,
                    href,
                    resolved_link_color,
                )
            continue
        if token.type == "hardbreak":
            add_linebreak()
            continue
        if token.type == "html_inline":
            if _validate_inline_html(token.content):
                add_linebreak()
            continue
        if token.type == "image":
            raise ValueError("inline image mixed with text is not supported")
        if token.type == "code_inline":
            _add_run(
                paragraph,
                token.content,
                mono_font,
                size_pt,
                color,
                state,
                href,
                resolved_link_color,
            )
            continue
        if token.type in ("text",):
            _add_run(
                paragraph,
                token.content,
                font,
                size_pt,
                color,
                state,
                href,
                resolved_link_color,
            )


def _set_bullet_glyph(
    paragraph: Any,
    glyph: str,
    color: str,
    level: int,
    font: str,
) -> None:
    """Give a textbox paragraph a real DrawingML bullet with a hanging indent."""
    hang = 0.28
    ppr = paragraph._p.get_or_add_pPr()
    ppr.set("marL", str(int((hang + level * hang) * 914400)))
    ppr.set("indent", str(int(-hang * 914400)))
    bu_clr = etree.SubElement(ppr, qn("a:buClr"))
    etree.SubElement(bu_clr, qn("a:srgbClr"), {"val": color})
    etree.SubElement(ppr, qn("a:buFont"), {"typeface": font})
    etree.SubElement(ppr, qn("a:buChar"), {"char": glyph})


# ============================================================
# Slide builder
# ============================================================


class PptxBuilder:
    def __init__(self, opts: Options) -> None:
        self.opts = opts
        self.theme = opts.theme
        self.palette = opts.palette
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W)
        self.prs.slide_height = Inches(SLIDE_H)
        self.blank = self.prs.slide_layouts[6]
        self.warnings: list[str] = []
        self.layouts_used: dict[str, int] = {}
        self.section_index = 0
        self.font_paths: dict[str, FontMetricSource] = {
            choice.role: FontMetricPlan(
                regular=(
                    (choice.metrics_path, choice.metrics_face_index)
                    if choice.metrics_path
                    else None
                ),
                styles=dict(choice.metrics_styles or {}),
            )
            for choice in opts.font_plan
        }

        for choice in opts.font_plan:
            if choice.substituted:
                self.warnings.append(
                    f"font '{choice.requested}' unavailable for role='{choice.role}', "
                    f"used '{choice.resolved}'"
                )
            elif not choice.verified:
                self.warnings.append(
                    f"font availability could not be verified for role='{choice.role}': "
                    f"'{choice.requested}'"
                )

    # ---- shared primitives ----

    def _new_slide(self, layout: str) -> Any:
        self.layouts_used[layout] = self.layouts_used.get(layout, 0) + 1
        return self.prs.slides.add_slide(self.blank)

    def _rect(self, slide: Any, x: float, y: float, w: float, h: float, color: str) -> Any:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color)
        shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def _textbox(self, slide: Any, x: float, y: float, w: float, h: float) -> Any:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.word_wrap = True
        frame.auto_size = MSO_AUTO_SIZE.NONE
        frame.margin_left = Inches(0)
        frame.margin_right = Inches(0)
        frame.margin_top = Inches(0)
        frame.margin_bottom = Inches(0)
        return box

    def _metrics_path_for_font(self, font: str) -> FontMetricSource:
        if font == self.theme.heading_font:
            return self.font_paths.get("heading")
        if font == self.theme.mono_font:
            return self.font_paths.get("mono")
        return self.font_paths.get("body")

    def _assert_text_fits(
        self,
        text: str,
        width_in: float,
        height_in: float,
        font_pt: float,
        font: str,
        label: str,
        bold: bool = False,
        italic: bool = False,
    ) -> None:
        lines = _estimate_lines(
            text,
            width_in,
            font_pt,
            self.opts.line_break_mode,
            self._metrics_path_for_font(font),
            self.font_paths.get("mono"),
            bold=bold,
            italic=italic,
        )
        estimated_height = lines * font_pt * 1.15 / POINTS_PER_INCH
        if estimated_height > height_in:
            raise ValueError(f"{label} text does not fit its fixed layout box")

    def _text(
        self,
        slide: Any,
        x: float,
        y: float,
        w: float,
        h: float,
        text: str,
        font: str,
        size_pt: float,
        color: str,
        bold: bool = False,
        italic: bool = False,
        align: Any = PP_ALIGN.LEFT,
        anchor: Any = MSO_ANCHOR.TOP,
        link_color: str | None = None,
    ) -> Any:
        self._assert_text_fits(
            text,
            w,
            h,
            size_pt,
            font,
            "slide",
            bold=bold,
            italic=italic,
        )
        box = self._textbox(slide, x, y, w, h)
        frame = box.text_frame
        frame.vertical_anchor = anchor
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        _add_md_runs(
            paragraph,
            text,
            font,
            size_pt,
            color,
            self.theme.mono_font,
            self.opts.line_break_mode,
            bold,
            italic,
            link_color=link_color or self.palette.link,
        )
        return box

    def _footer(self, slide: Any, number: int) -> None:
        if self.opts.footer_text:
            self._text(
                slide,
                MARGIN,
                7.08,
                6.0,
                0.3,
                self.opts.footer_text,
                self.theme.body_font,
                10,
                self.palette.text_muted,
            )
        if self.opts.slide_numbers:
            self._text(
                slide,
                SLIDE_W - MARGIN - 0.6,
                7.08,
                0.6,
                0.3,
                str(number),
                self.theme.body_font,
                10,
                self.palette.text_muted,
                align=PP_ALIGN.RIGHT,
            )

    def _content_header(self, slide: Any, title: str | None, kicker: str | None) -> None:
        if kicker:
            self._text(
                slide,
                MARGIN,
                0.52,
                CONTENT_W,
                0.32,
                kicker.upper(),
                self.theme.body_font,
                11,
                self.palette.accent_text,
                bold=True,
            )
        if title:
            self._text(
                slide,
                MARGIN,
                0.85,
                CONTENT_W,
                0.75,
                title,
                self.theme.heading_font,
                24,
                self.palette.text,
                bold=True,
            )
            self._rect(slide, MARGIN, 1.68, 1.1, 0.045, self.palette.accent)

    def _bullets_into(
        self, frame: Any, bullets: list[dict[str, Any]], sizes: tuple[float, float, float]
    ) -> None:
        glyphs = [
            ("•", self.palette.accent),
            ("–", self.palette.text_muted),
            ("·", self.palette.text_muted),
        ]
        first = True
        for bullet in bullets:
            paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
            first = False
            level = int(bullet["level"])
            size = sizes[level]
            paragraph.line_spacing = Pt(size * BULLET_LINE_HEIGHT_FACTOR)
            paragraph.space_after = Pt(max(4.0, size * BULLET_GAP_FACTOR))
            glyph, glyph_color = glyphs[level]
            _set_bullet_glyph(
                paragraph,
                glyph,
                glyph_color,
                level,
                self.theme.body_font,
            )
            color = self.palette.text if level == 0 else self.palette.text_muted
            _add_md_runs(
                paragraph,
                str(bullet["text"]),
                self.theme.body_font,
                size,
                color,
                self.theme.mono_font,
                self.opts.line_break_mode,
                bold=bool(bullet.get("bold")),
                link_color=self.palette.link,
            )

    def _notes(self, slide: Any, notes: str | None) -> None:
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    # ---- pages ----

    def _content_pages_for_spec(self, spec: dict[str, Any]) -> list[dict[str, Any]]:
        title = spec.get("title")
        kicker = spec.get("kicker")

        body_parts = [str(p) for p in (spec.get("body_parts") or []) if p is not None]
        if not body_parts and spec.get("body"):
            body_parts = [str(spec["body"])]

        if spec.get("code_lines") is not None:
            code_parts = ["\n".join(str(line) for line in spec.get("code_lines") or [])]
        else:
            code_parts = [str(code) for code in (spec.get("code") or spec.get("code_parts") or [])]

        pages = _paginate_content(
            title=title,
            kicker=kicker,
            body_parts=body_parts,
            bullets=[_normalize_bullet(item) for item in (spec.get("bullets") or [])],
            code_parts=code_parts,
            line_break_mode=self.opts.line_break_mode,
            body_font_path=self.font_paths.get("body"),
            mono_font_path=self.font_paths.get("mono"),
        )
        if spec.get("notes") and pages:
            pages[0]["notes"] = spec["notes"]
        return pages

    def _table_pages_for_spec(self, spec: dict[str, Any]) -> list[dict[str, Any]]:
        table = spec.get("table") or {}
        title = spec.get("title")
        # Keep compatibility with already segmented table specs.
        if isinstance(table, dict) and table.get("_prepaginated"):
            return [spec]
        body_parts = [str(part) for part in (spec.get("body_parts") or [])]
        if not body_parts and spec.get("body"):
            body_parts = _body_parts(spec.get("body"))
        pages = _paginate_table(
            title=title,
            kicker=spec.get("kicker"),
            body_parts=body_parts,
            table=table if isinstance(table, dict) else {"headers": [], "rows": []},
            line_break_mode=self.opts.line_break_mode,
            body_font_path=self.font_paths.get("body"),
            mono_font_path=self.font_paths.get("mono"),
        )
        if spec.get("notes") and pages:
            pages[0]["notes"] = spec["notes"]
        return pages

    def _pages_from(self, spec: dict[str, Any]) -> list[dict[str, Any]]:
        layout = spec["layout"]
        if layout == "content":
            return self._content_pages_for_spec(spec)
        if layout == "table":
            return self._table_pages_for_spec(spec)
        return [spec]

    def build(self, spec: dict[str, Any], number: int) -> int:
        pages = self._pages_from(spec)
        for page in pages:
            layout = page["layout"]
            builder = getattr(self, f"_build_{layout}")
            slide = builder(page)
            if layout in ("content", "two_column", "table", "quote", "image"):
                self._footer(slide, number)
            self._notes(slide, page.get("notes"))
            number += 1
        return len(pages)

    # ---- layouts ----

    def _build_title(self, spec: dict[str, Any]) -> Any:
        slide = self._new_slide("title")
        self._rect(slide, MARGIN, 2.35, 1.6, 0.055, self.palette.accent)
        self._text(
            slide,
            MARGIN,
            2.55,
            11.0,
            1.7,
            spec.get("title") or "Presentation",
            self.theme.heading_font,
            40,
            self.palette.text,
            bold=True,
        )
        subtitle = spec.get("subtitle") or self.opts.subtitle
        if subtitle:
            self._text(
                slide,
                MARGIN,
                4.2,
                10.5,
                0.9,
                str(subtitle),
                self.theme.body_font,
                18,
                self.palette.text_muted,
            )
        meta = "  ·  ".join(
            p
            for p in (spec.get("author") or self.opts.author, spec.get("date") or self.opts.date)
            if p
        )
        if meta:
            self._text(
                slide,
                MARGIN,
                6.35,
                10.0,
                0.45,
                meta,
                self.theme.body_font,
                12,
                self.palette.text_muted,
            )
        self._rect(slide, 11.6, 5.9, 1.05, 1.05, self.palette.accent_soft)
        self._rect(slide, 12.15, 6.45, 0.5, 0.5, self.palette.accent)
        return slide

    def _build_section(self, spec: dict[str, Any]) -> Any:
        slide = self._new_slide("section")
        self.section_index += 1
        self._rect(slide, 0, 0, SLIDE_W, SLIDE_H, self.palette.accent)
        self._text(
            slide,
            10.4,
            0.45,
            2.4,
            1.7,
            f"{self.section_index:02d}",
            self.theme.heading_font,
            96,
            _mix(self.palette.accent, self.palette.on_accent, 0.25),
            bold=True,
            align=PP_ALIGN.RIGHT,
        )
        self._rect(slide, MARGIN, 2.95, 1.2, 0.055, self.palette.on_accent)
        self._text(
            slide,
            MARGIN,
            3.15,
            11.9,
            1.6,
            spec.get("title") or "",
            self.theme.heading_font,
            32,
            self.palette.on_accent,
            bold=True,
            link_color=self.palette.link_on_accent,
        )
        return slide

    def _build_content(self, spec: dict[str, Any]) -> Any:
        slide = self._new_slide("content")
        self._content_header(slide, spec.get("title"), spec.get("kicker"))
        y = CONTENT_TOP

        body_parts = [str(p) for p in (spec.get("body_parts") or [])]
        if body_parts:
            y = self._body_text(slide, y, body_parts)

        code_lines = [str(line) for line in (spec.get("code_lines") or [])]
        if code_lines:
            y = self._code_panel(slide, y, code_lines)

        bullets = [dict(item) for item in (spec.get("bullets") or [])]
        if bullets:
            sizes = tuple(spec.get("_bullet_sizes") or BULLET_FONT_TIERS[0])
            height = sum(
                _bullet_height(
                    bullet,
                    sizes,
                    self.opts.line_break_mode,
                    self.font_paths.get("body"),
                    self.font_paths.get("mono"),
                )
                for bullet in bullets
            )
            if y + height > CONTENT_BOTTOM + 1e-6:
                raise ValueError("bullet content overflows the planned content area")
            box = self._textbox(slide, MARGIN, y, CONTENT_W, max(0.5, height))
            self._bullets_into(box.text_frame, bullets, sizes)
        return slide

    def _body_text(
        self,
        slide: Any,
        y: float,
        body_parts: list[str],
        x: float = MARGIN,
        width: float = CONTENT_W,
        font_pt: float = BODY_FONT_PT,
        line_height_in: float = BODY_LINE_HEIGHT_IN,
    ) -> float:
        height = sum(
            _body_paragraph_height(
                part,
                self.opts.line_break_mode,
                self.font_paths.get("body"),
                self.font_paths.get("mono"),
                width_in=width,
                font_pt=font_pt,
                line_height_in=line_height_in,
            )
            for part in body_parts
        )
        if y + height > CONTENT_BOTTOM + 1e-6:
            raise ValueError("body content overflows the planned content area")
        box = self._textbox(slide, x, y, width, height)
        first = True
        for part in body_parts:
            paragraph = box.text_frame.paragraphs[0] if first else box.text_frame.add_paragraph()
            first = False
            paragraph.line_spacing = Pt(line_height_in * POINTS_PER_INCH)
            paragraph.space_after = Pt(PARAGRAPH_GAP_IN * POINTS_PER_INCH)
            _add_md_runs(
                paragraph,
                str(part),
                self.theme.body_font,
                font_pt,
                self.palette.text,
                self.theme.mono_font,
                self.opts.line_break_mode,
                link_color=self.palette.link,
            )
        return y + height

    def _code_panel(self, slide: Any, y: float, code_lines: list[str]) -> float:
        if not code_lines:
            return y
        content_height = sum(
            _code_source_line_height(line, self.font_paths.get("mono")) for line in code_lines
        )
        height = content_height + CODE_PANEL_PADDING_IN
        if y + height > CONTENT_BOTTOM + 1e-6:
            raise ValueError("code content overflows the planned content area")
        self._rect(slide, MARGIN, y, CONTENT_W, height, self.palette.code_background)
        box = self._textbox(
            slide,
            MARGIN + 0.15,
            y + CODE_PANEL_PADDING_IN / 2,
            CONTENT_W - 0.3,
            content_height,
        )
        first = True
        for line in code_lines:
            paragraph = box.text_frame.paragraphs[0] if first else box.text_frame.add_paragraph()
            first = False
            paragraph.line_spacing = Pt(CODE_LINE_HEIGHT_IN * POINTS_PER_INCH)
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(0)
            run = paragraph.add_run()
            run.text = str(line) if line else " "
            run.font.name = self.theme.mono_font
            run.font.size = Pt(CODE_FONT_PT)
            run.font.color.rgb = RGBColor.from_string(self.palette.text)
        return y + height

    def _build_two_column(self, spec: dict[str, Any]) -> Any:
        slide = self._new_slide("two_column")
        self._content_header(slide, spec.get("title"), spec.get("kicker"))
        column_w = (CONTENT_W - 0.6) / 2
        body_font_pt = 13.0
        body_line_height = body_font_pt * BULLET_LINE_HEIGHT_FACTOR / POINTS_PER_INCH
        for index, side in enumerate(("left", "right")):
            column = spec.get(side) or {}
            x = MARGIN + index * (column_w + 0.6)
            y = CONTENT_TOP
            if column.get("heading"):
                self._text(
                    slide,
                    x,
                    y,
                    column_w,
                    0.4,
                    str(column["heading"]),
                    self.theme.body_font,
                    15,
                    self.palette.accent_text,
                    bold=True,
                )
                y += 0.5
            column_body_parts = [str(part) for part in (column.get("body_parts") or [])]
            if not column_body_parts and column.get("body"):
                column_body_parts = _body_parts(column.get("body"))
            if column_body_parts:
                try:
                    y = self._body_text(
                        slide,
                        y,
                        column_body_parts,
                        x=x,
                        width=column_w,
                        font_pt=body_font_pt,
                        line_height_in=body_line_height,
                    )
                except ValueError as exc:
                    raise ValueError(
                        f"{side} column content does not fit its fixed layout box"
                    ) from exc
            bullets = [_normalize_bullet(item) for item in column.get("bullets") or []]
            if bullets:
                sizes = BULLET_FONT_TIERS[-1]
                bullet_height = sum(
                    _bullet_height(
                        bullet,
                        sizes,
                        self.opts.line_break_mode,
                        self.font_paths.get("body"),
                        self.font_paths.get("mono"),
                        width_in=column_w,
                    )
                    for bullet in bullets
                )
                if y + bullet_height > CONTENT_BOTTOM + 1e-6:
                    raise ValueError(f"{side} column content does not fit its fixed layout box")
                box = self._textbox(slide, x, y, column_w, bullet_height)
                self._bullets_into(box.text_frame, bullets, sizes)
        return slide

    def _build_table(self, spec: dict[str, Any]) -> Any:
        slide = self._new_slide("table")
        self._content_header(slide, spec.get("title"), spec.get("kicker"))
        y = CONTENT_TOP
        body_parts = [str(part) for part in (spec.get("body_parts") or [])]
        if not body_parts and spec.get("body"):
            body_parts = _body_parts(spec.get("body"))
        if body_parts:
            y = self._body_text(slide, y, body_parts)

        table_spec = spec.get("table") or {}
        headers = [str(h) for h in (table_spec.get("headers") or [])]
        rows = [[str(c) for c in row] for row in (table_spec.get("rows") or [])]
        cols = max(len(headers), max((len(row) for row in rows), default=0), 1)
        column_width = CONTENT_W / cols
        header_height = float(
            table_spec.get("header_height")
            or _table_row_height(
                headers,
                column_width,
                self.opts.line_break_mode,
                self.font_paths.get("body"),
                self.font_paths.get("mono"),
                font_pt=TABLE_HEADER_FONT_PT,
                minimum=TABLE_HEADER_MIN_HEIGHT_IN,
                bold=True,
            )
        )
        row_heights = [float(height) for height in table_spec.get("row_heights") or []]
        if len(row_heights) != len(rows):
            row_heights = [
                _table_row_height(
                    row,
                    column_width,
                    self.opts.line_break_mode,
                    self.font_paths.get("body"),
                    self.font_paths.get("mono"),
                )
                for row in rows
            ]
        table_height = header_height + sum(row_heights)
        if y + table_height > CONTENT_BOTTOM + 1e-6:
            raise ValueError("table content overflows the planned content area")
        shape = slide.shapes.add_table(
            len(rows) + 1,
            cols,
            Inches(MARGIN),
            Inches(y),
            Inches(CONTENT_W),
            Inches(table_height),
        )
        self._style_table(
            shape.table,
            headers,
            rows,
            cols,
            header_height,
            row_heights,
        )
        return slide

    def _style_table(
        self,
        table: Any,
        headers: list[str],
        rows: list[list[str]],
        cols: int,
        header_height: float,
        row_heights: list[float],
    ) -> None:
        table.first_row = False
        table.horz_banding = False
        table.rows[0].height = Inches(header_height)
        for index, height in enumerate(row_heights):
            table.rows[index + 1].height = Inches(height)
        for c in range(cols):
            header = headers[c] if c < len(headers) else ""
            self._style_cell(
                table.cell(0, c),
                header,
                self.palette.accent,
                self.palette.on_accent,
                TABLE_HEADER_FONT_PT,
                bold=True,
                link_color=self.palette.link_on_accent,
            )
            for r, row in enumerate(rows):
                fill = BAND_FILL if r % 2 == 1 else "FFFFFF"
                value = row[c] if c < len(row) else ""
                self._style_cell(
                    table.cell(r + 1, c),
                    value,
                    fill,
                    self.palette.text,
                    TABLE_BODY_FONT_PT,
                    link_color=self.palette.link,
                )

    def _style_cell(
        self,
        cell: Any,
        text: str,
        fill: str,
        color: str,
        size_pt: float,
        bold: bool = False,
        link_color: str | None = None,
    ) -> None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(fill)
        cell.margin_left = Inches(0.12)
        cell.margin_right = Inches(0.12)
        cell.margin_top = Inches(TABLE_CELL_VERTICAL_PADDING_IN / 2)
        cell.margin_bottom = Inches(TABLE_CELL_VERTICAL_PADDING_IN / 2)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = cell.text_frame.paragraphs[0]
        _add_md_runs(
            paragraph,
            text,
            self.theme.body_font,
            size_pt,
            color,
            self.theme.mono_font,
            self.opts.line_break_mode,
            bold=bold,
            link_color=link_color or self.palette.link,
        )

    def _build_quote(self, spec: dict[str, Any]) -> Any:
        slide = self._new_slide("quote")
        has_header = bool(spec.get("title") or spec.get("kicker"))
        if has_header:
            self._content_header(slide, spec.get("title"), spec.get("kicker"))
        glyph_y, quote_y, attr_y = (1.85, 3.0, 5.6) if has_header else (0.55, 2.45, 5.35)
        self._text(
            slide,
            0.9,
            glyph_y,
            2.2,
            2.0,
            "\u201c",
            self.theme.heading_font,
            120,
            self.palette.accent_soft,
            bold=True,
        )
        self._text(
            slide,
            2.17,
            quote_y,
            9.0,
            2.4,
            str(spec.get("quote") or ""),
            self.theme.heading_font,
            24,
            self.palette.text,
            italic=True,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        if spec.get("attribution"):
            self._text(
                slide,
                2.17,
                attr_y,
                9.0,
                0.5,
                f"\u2014 {spec['attribution']}",
                self.theme.body_font,
                14,
                self.palette.text_muted,
                align=PP_ALIGN.CENTER,
            )
        return slide

    def _build_image(self, spec: dict[str, Any]) -> Any:
        slide = self._new_slide("image")
        self._content_header(slide, spec.get("title"), spec.get("kicker"))
        src = str(spec.get("image_path") or "")
        path = src if os.path.isabs(src) else os.path.join(self.opts.project_root, src)
        caption = spec.get("caption")
        top = CONTENT_TOP if spec.get("title") else 0.9
        bottom = 6.45 if caption else CONTENT_BOTTOM

        if not os.path.isfile(path):
            self.warnings.append(f"image not found: {src}")
            self._text(
                slide,
                MARGIN,
                3.2,
                CONTENT_W,
                0.6,
                f"[image unavailable: {caption or src}]",
                self.theme.body_font,
                14,
                self.palette.text_muted,
                align=PP_ALIGN.CENTER,
            )
            return slide

        image_source, width, height, ppi_x, ppi_y = self._prepare_image(
            path,
            CONTENT_W,
            bottom - top,
        )
        if ppi_x < MIN_EFFECTIVE_PPI:
            self.warnings.append(f"low effective horizontal PPI for image {src}: {ppi_x:.1f}")
        if ppi_y < MIN_EFFECTIVE_PPI:
            self.warnings.append(f"low effective vertical PPI for image {src}: {ppi_y:.1f}")

        slide.shapes.add_picture(
            image_source,
            Inches(MARGIN + (CONTENT_W - width) / 2),
            Inches(top + (bottom - top - height) / 2),
            Inches(width),
            Inches(height),
        )
        if caption:
            self._text(
                slide,
                MARGIN,
                6.55,
                CONTENT_W,
                0.4,
                str(caption),
                self.theme.body_font,
                11,
                self.palette.text_muted,
                align=PP_ALIGN.CENTER,
            )
        return slide

    def _prepare_image(
        self,
        path: str,
        max_w: float,
        max_h: float,
    ) -> tuple[str | io.BytesIO, float, float, float, float]:
        try:
            with Image.open(path) as image:
                image.load()
                orientation = image.getexif().get(274, 1)
                is_animated = bool(getattr(image, "is_animated", False))
                if is_animated and orientation not in (None, 1):
                    raise ValueError(
                        "animated images with EXIF orientation are not supported safely"
                    )
                normalized = ImageOps.exif_transpose(image)
                width_px, height_px = normalized.size
                if width_px <= 0 or height_px <= 0:
                    raise ValueError("image dimensions must be positive")
                image_source: str | io.BytesIO = path
                if orientation not in (None, 1):
                    payload = io.BytesIO()
                    normalized.save(payload, format="PNG")
                    payload.seek(0)
                    image_source = payload
        except ValueError:
            raise
        except Exception as exc:
            raise ValueError(f"corrupt image: {path}") from exc

        ratio = width_px / height_px
        width_rendered = min(max_w, max_h * ratio)
        height_rendered = width_rendered / ratio
        effective_ppi_x = width_px / width_rendered
        effective_ppi_y = height_px / height_rendered
        return image_source, width_rendered, height_rendered, effective_ppi_x, effective_ppi_y

    def _build_closing(self, spec: dict[str, Any]) -> Any:
        slide = self._new_slide("closing")
        self._rect(slide, 0, 0, SLIDE_W, SLIDE_H, self.palette.accent)
        self._text(
            slide,
            1.17,
            2.9,
            11.0,
            1.2,
            spec.get("title") or "Thank you",
            self.theme.heading_font,
            32,
            self.palette.on_accent,
            bold=True,
            align=PP_ALIGN.CENTER,
            link_color=self.palette.link_on_accent,
        )
        subtitle = spec.get("subtitle")
        if subtitle:
            self._text(
                slide,
                1.17,
                4.15,
                11.0,
                0.8,
                str(subtitle),
                self.theme.body_font,
                16,
                self.palette.on_accent,
                align=PP_ALIGN.CENTER,
                link_color=self.palette.link_on_accent,
            )
        return slide


# ============================================================
# Entry point
# ============================================================


def _load_slides(spec: dict[str, Any], opts: Options) -> tuple[list[dict[str, Any]], list[str]]:
    slides = spec.get("slides")
    markdown = spec.get("markdown")
    # Whitespace-only markdown counts as absent, matching the TS-side gate.
    has_markdown = bool(markdown and str(markdown).strip())
    if slides and has_markdown:
        raise ValueError("provide exactly one of 'slides' or 'markdown', not both")
    if slides:
        if not isinstance(slides, list):
            raise ValueError("'slides' must be a non-empty array")
        return [normalize_slide(s) for s in slides], []
    if has_markdown:
        result, warnings = markdown_to_slides(str(markdown), opts)
        if not result:
            raise ValueError("markdown produced no slides")
        return result, warnings
    raise ValueError("spec requires 'slides' or non-empty 'markdown'")


def _parse_package_xml(payload: bytes, part_name: str) -> None:
    parser = etree.XMLParser(
        resolve_entities=False,
        no_network=True,
        load_dtd=False,
        recover=False,
        huge_tree=False,
    )
    try:
        tree = etree.parse(io.BytesIO(payload), parser)
    except (etree.XMLSyntaxError, ValueError) as exc:
        raise ValueError(f"PPTX XML parse failure in {part_name}: {exc}") from exc
    if tree.docinfo.doctype:
        raise ValueError(f"PPTX XML contains a forbidden DTD/entity declaration in {part_name}")


def _validate_pptx(
    path: str | Path,
    expected_slide_count: int | None = None,
) -> dict[str, Any]:
    package_path = Path(path)
    xml_parts: list[str] = []
    try:
        with zipfile.ZipFile(package_path, "r") as archive:
            names = archive.namelist()
            if len(names) != len(set(names)):
                duplicates = sorted({name for name in names if names.count(name) > 1})
                raise ValueError(f"PPTX contains duplicate part names: {duplicates}")
            name_set = set(names)
            missing = _REQUIRED_PPTX_PARTS - name_set
            if missing:
                raise ValueError(f"PPTX is missing required parts: {sorted(missing)}")
            broken = archive.testzip()
            if broken is not None:
                raise ValueError(f"PPTX archive corrupt: {broken!r}")
            xml_parts = [name for name in names if name.lower().endswith((".xml", ".rels"))]
            for name in xml_parts:
                _parse_package_xml(archive.read(name), name)
    except zipfile.BadZipFile as exc:
        raise ValueError(f"invalid generated PPTX: {exc}") from exc

    reopened = Presentation(str(package_path))
    slide_count = len(reopened.slides)
    if expected_slide_count is not None and slide_count != expected_slide_count:
        raise ValueError(
            f"PPTX slide count mismatch: expected {expected_slide_count}, got {slide_count}"
        )
    return {
        "package_valid": True,
        "reopen_valid": True,
        "slide_count": slide_count,
        "required_parts": sorted(_REQUIRED_PPTX_PARTS),
        "xml_parts_checked": len(xml_parts),
        "openxml_schema_validation": "not_performed",
    }


def _atomic_save(
    presentation: PresentationType,
    output_path: str | Path,
    staging_path: str | Path | None,
) -> tuple[str, dict[str, Any]]:
    target = Path(output_path).expanduser().resolve()
    target.parent.mkdir(parents=True, exist_ok=True)

    if staging_path is not None:
        staging = Path(staging_path).expanduser().resolve()
        if staging == target:
            raise ValueError("staging_path must be distinct from output_path")
        if staging.parent != target.parent:
            raise ValueError("staging_path must be in the same directory as output_path")
    else:
        descriptor, staged = tempfile.mkstemp(
            prefix=f".{target.stem}.",
            suffix=".tmp.pptx",
            dir=target.parent,
        )
        os.close(descriptor)
        staging = Path(staged)

    try:
        presentation.save(str(staging))
        with staging.open("rb") as handle:
            os.fsync(handle.fileno())
        validation = _validate_pptx(
            staging,
            expected_slide_count=len(presentation.slides),
        )
        os.replace(staging, target)
        return str(target), validation
    finally:
        staging.unlink(missing_ok=True)


def generate(spec: dict[str, Any]) -> dict[str, Any]:
    opts = parse_options(spec)
    slides, md_warnings = _load_slides(spec, opts)

    builder = PptxBuilder(opts)
    builder.warnings.extend(md_warnings)

    number = 1
    for slide_spec in slides:
        number += builder.build(slide_spec, number)

    target_path = opts.output_path
    resolved_target, validation = _atomic_save(builder.prs, target_path, opts.staging_path)
    font_records = [asdict(choice) for choice in opts.font_plan]
    continued_tables = max(
        0,
        builder.layouts_used.get("table", 0)
        - sum(1 for slide in slides if slide.get("layout") == "table"),
    )
    continued_content = max(
        0,
        builder.layouts_used.get("content", 0)
        - sum(1 for slide in slides if slide.get("layout") == "content"),
    )

    return {
        "path": os.path.abspath(resolved_target),
        "slide_count": len(builder.prs.slides),
        "layouts_used": builder.layouts_used,
        "theme": opts.theme_name,
        "warnings": builder.warnings,
        "validation": validation,
        "resolved_palette": asdict(opts.palette),
        "fonts": font_records,
        "font_plan": font_records,
        "line_break_mode": opts.line_break_mode,
        "normalization": {
            "line_break_mode": opts.line_break_mode,
            "split_slides": max(0, len(builder.prs.slides) - len(slides)),
            "continued_tables": continued_tables,
            "continued_content_slides": continued_content,
        },
    }


def main() -> None:
    spec = json.loads(sys.stdin.read())
    result = generate(spec)
    print(json.dumps(result))


if __name__ == "__main__":
    try:
        main()
    except Exception as exc:  # noqa: BLE001 — stable stderr contract for the TS caller
        print(f"{type(exc).__name__}: {exc}", file=sys.stderr)
        sys.exit(1)
